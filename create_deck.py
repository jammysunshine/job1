#!/usr/bin/env python3
"""
MYTAR — Full Strawman Investor Deck (Expanded Version)
Fully designed VC-style PPTX with visuals, diagrams, and icons.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Colour Palette ──────────────────────────────────────────────────────────
BG_DARK       = RGBColor(0x0B, 0x0E, 0x1A)   # near-black navy
BG_MID        = RGBColor(0x11, 0x16, 0x28)   # slightly lighter
ACCENT_TEAL   = RGBColor(0x33, 0xE0, 0xCC)   # bright teal
ACCENT_BLUE   = RGBColor(0x29, 0x80, 0xB9)   # medium blue
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)   # purple
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY    = RGBColor(0xBB, 0xBF, 0xD0)
DIM_GREY      = RGBColor(0x6B, 0x72, 0x8A)
ACCENT_GREEN  = RGBColor(0x10, 0xB9, 0x81)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_RED    = RGBColor(0xEF, 0x44, 0x44)

TOTAL = 44

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── Helper Functions ────────────────────────────────────────────────────────

def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    return slide


def bg(slide, color=BG_DARK):
    """Fill slide background with a solid colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, fill_color=ACCENT_TEAL, border=None, radius=None):
    """Add a rectangle shape."""
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = radius
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def circle(slide, left, top, size, fill_color=ACCENT_TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def text_box(slide, left, top, width, height, text="", font_size=14,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Calibri", anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
   	 p.line_spacing = Pt(font_size * line_spacing)
    tf.paragraphs[0].font.name = font_name
    # set anchor
    txBox.text_frame.paragraphs[0].font.name = font_name
    try:
   	 txBox.text_frame.auto_size = None
    except:
   	 pass
    return txBox


def multi_text(slide, left, top, width, height, lines, font_name="Calibri",
               anchor=MSO_ANCHOR.TOP):
    """lines = [(text, font_size, color, bold, alignment), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, (text, font_size, color, bold, alignment) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
        p.space_before = Pt(2)
    return txBox


def accent_line(slide, left, top, width=Inches(0.6), height=Pt(4), color=ACCENT_TEAL):
    return rect(slide, left, top, width, height, fill_color=color)


def page_number(slide, num=None, total=TOTAL):
    if num is None:
        num = len(prs.slides)
    text_box(slide, W - Inches(1.5), H - Inches(0.5), Inches(1.2), Inches(0.4),
             f"{num:02d} / {total:02d}", font_size=9, color=DIM_GREY,
             alignment=PP_ALIGN.RIGHT)


def add_icon(slide, icon_char, left, top, size=Inches(0.5), color=ACCENT_TEAL,
             bg_color=None):
    """Add a circular icon with a unicode character."""
    if bg_color:
        c = circle(slide, left, top, size, fill_color=bg_color)
    else:
        c = circle(slide, left, top, size, fill_color=RGBColor(0x1A, 0x22, 0x3A))
    cx = c.left + c.width // 2
    cy = c.top + c.height // 2
    tbx = text_box(slide, left, top, size, size, icon_char,
                   font_size=int(size / Pt(1) * 0.45), color=color,
                   bold=False, alignment=PP_ALIGN.CENTER, font_name="Segoe UI Symbol")
    tbx.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return c


def add_section_header(slide, number, title, subtitle=None, accent=True):
    """Standard section header treatment."""
    bg(slide, BG_DARK)
    # decorative top bar
    rect(slide, 0, 0, W, Pt(4), fill_color=ACCENT_TEAL)
    # left accent stripe
    rect(slide, Inches(0.6), Inches(0.8), Pt(4), Inches(1.2), fill_color=ACCENT_TEAL)
    # section number
    text_box(slide, Inches(1.0), Inches(0.8), Inches(2), Inches(0.6),
             f"0{number}" if number < 10 else str(number),
             font_size=42, color=DIM_GREY, bold=True, font_name="Calibri")
    # title
    text_box(slide, Inches(1.0), Inches(1.6), W - Inches(3), Inches(0.8),
             title, font_size=36, color=WHITE, bold=True, font_name="Calibri")
    if subtitle:
        text_box(slide, Inches(1.0), Inches(2.5), W - Inches(3), Inches(0.5),
                 subtitle, font_size=16, color=LIGHT_GREY, font_name="Calibri")
    # bottom gradient bar
    rect(slide, 0, H - Inches(0.06), W, Pt(3), fill_color=ACCENT_TEAL)


def bullet_block(slide, left, top, width, height, bullets, font_size=14,
                 color=LIGHT_GREY, bullet_char="\u25B8", font_name="Calibri"):
    """Add a block of bullet points."""
    lines = []
    for b in bullets:
        if isinstance(b, tuple):
            txt, sz, clr, bld = b
            lines.append((f"{bullet_char} {txt}", sz, clr, bld, PP_ALIGN.LEFT))
        else:
            lines.append((f"{bullet_char} {b}", font_size, color, False, PP_ALIGN.LEFT))
    return multi_text(slide, left, top, width, height, lines, font_name=font_name)


# ── SLIDE 1 ── Title Slide ────────────────────────────────────────────────────
def slide_01():
    s = add_blank_slide()
    bg(s, BG_DARK)
    # large background accent circle
    circ = circle(s, W - Inches(2.5), -Inches(1.5), Inches(5), fill_color=RGBColor(0x12, 0x1C, 0x38))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(0x12, 0x1C, 0x38)

    circ2 = circle(s, -Inches(1), H - Inches(3), Inches(4), fill_color=RGBColor(0x0E, 0x16, 0x2E))

    # top bar
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_TEAL)

    # Brand
    text_box(s, Inches(0.8), Inches(0.6), Inches(4), Inches(0.5),
             "MYTAR", font_size=20, color=ACCENT_TEAL, bold=True, font_name="Calibri")

    # Main title
    text_box(s, Inches(0.8), Inches(2.0), Inches(8), Inches(1.2),
             "My Avatar Platform", font_size=56, color=WHITE, bold=True, font_name="Calibri")

    # Tagline
    text_box(s, Inches(0.8), Inches(3.3), Inches(8), Inches(0.8),
             "The Ownership Layer for Human Identity in the AI Era",
             font_size=22, color=LIGHT_GREY, font_name="Calibri")

    # Accent line
    rect(s, Inches(0.8), Inches(4.3), Inches(1.5), Pt(4), fill_color=ACCENT_TEAL)

    # Subtitle
    text_box(s, Inches(0.8), Inches(4.6), Inches(7), Inches(0.6),
             "Investor Presentation  \u2022  Confidential",
             font_size=14, color=DIM_GREY, font_name="Calibri")

    # Decorative dots
    circle(s, W - Inches(2.5), Inches(1.0), Inches(0.08), fill_color=ACCENT_TEAL)
    circle(s, W - Inches(2.5), Inches(1.3), Inches(0.08), fill_color=DIM_GREY)
    circle(s, W - Inches(2.5), Inches(1.6), Inches(0.08), fill_color=DIM_GREY)

    # Right side decorative element
    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_TEAL)

    page_number(s)


# ── SLIDE 2 ── Founding Vision ────────────────────────────────────────────────
def slide_02():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_PURPLE)

    # Vision quote style
    text_box(s, Inches(0.8), Inches(1.5), Inches(1.5), Inches(0.5),
             "VISION", font_size=12, color=ACCENT_PURPLE, bold=True)
    rect(s, Inches(0.8), Inches(2.1), Inches(0.8), Pt(3), fill_color=ACCENT_PURPLE)

    main_lines = [
        ("The world is transitioning into a phase where:", 22, WHITE, False, PP_ALIGN.LEFT),
        ("", 10, WHITE, False, PP_ALIGN.LEFT),
        ("\u25B8  Humans can be perfectly simulated", 18, LIGHT_GREY, False, PP_ALIGN.LEFT),
        ("\u25B8  Voice, face, and behavior can be cloned", 18, LIGHT_GREY, False, PP_ALIGN.LEFT),
        ("\u25B8  Digital presence is no longer inherently trustworthy", 18, LIGHT_GREY, False, PP_ALIGN.LEFT),
        ("", 10, WHITE, False, PP_ALIGN.LEFT),
        ("In this world, identity becomes the most valuable asset.", 20, ACCENT_TEAL, False, PP_ALIGN.LEFT),
        ("", 10, WHITE, False, PP_ALIGN.LEFT),
        ("MYTAR emerges from a simple belief:", 18, WHITE, False, PP_ALIGN.LEFT),
        ("", 6, WHITE, False, PP_ALIGN.LEFT),
        ("Every human should own a cryptographically verifiable", 26, WHITE, True, PP_ALIGN.LEFT),
        ("version of themselves.", 26, ACCENT_TEAL, True, PP_ALIGN.LEFT),
    ]
    multi_text(s, Inches(0.8), Inches(2.6), Inches(11), Inches(4.5), main_lines)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_PURPLE)
    page_number(s)


# ── SLIDE 3 ── Core Thesis ────────────────────────────────────────────────────
def slide_03():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_TEAL)

    text_box(s, Inches(0.8), Inches(1.5), Inches(1.5), Inches(0.5),
             "THESIS", font_size=12, color=ACCENT_TEAL, bold=True)
    rect(s, Inches(0.8), Inches(2.1), Inches(0.8), Pt(3), fill_color=ACCENT_TEAL)

    main_lines = [
        ("\"You are not just a user on the internet.", 28, WHITE, False, PP_ALIGN.LEFT),
        ("You are a verifiable digital entity with", 28, WHITE, False, PP_ALIGN.LEFT),
        ("ownership, rights, and representation.\"", 28, ACCENT_TEAL, True, PP_ALIGN.LEFT),
        ("", 10, WHITE, False, PP_ALIGN.LEFT),
        ("This includes:", 16, LIGHT_GREY, True, PP_ALIGN.LEFT),
        ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ]
    multi_text(s, Inches(0.8), Inches(2.6), Inches(11), Inches(3), main_lines)

    # 4 feature boxes
    features = [
        ("\U0001F511", "Blockchain-registered\nidentity", "Cryptographic proof\nof personhood"),
        ("\U0001F3F7", "NFT Identity\nCertificate", "Ownership deed for\nyour digital self"),
        ("\U0001F9E0", "Evolving AI\nTwin", "Continuous learning\nrepresentation"),
        ("\u270D\ufe0f", "Content Signature\nSystem", "Authenticity proof\nfor every output"),
    ]
    for i, (icon, title, desc) in enumerate(features):
        x = Inches(0.8 + i * 3.1)
        y = Inches(4.5)
        # box background
        rect(s, x, y, Inches(2.8), Inches(2.2), fill_color=BG_MID, radius=0.05)
        # accent top
        rect(s, x, y, Inches(2.8), Pt(3), fill_color=ACCENT_TEAL)
        # icon
        text_box(s, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5),
                 icon, font_size=24, color=ACCENT_TEAL, font_name="Segoe UI Symbol")
        # title
        text_box(s, x + Inches(0.2), y + Inches(0.8), Inches(2.4), Inches(0.6),
                 title, font_size=14, color=WHITE, bold=True)
        # desc
        text_box(s, x + Inches(0.2), y + Inches(1.4), Inches(2.4), Inches(0.6),
                 desc, font_size=11, color=LIGHT_GREY)

    text_box(s, Inches(0.8), Inches(6.9), Inches(8), Inches(0.4),
             "This shifts the internet from: anonymous \u2192 authenticated humans",
             font_size=13, color=DIM_GREY, font_name="Calibri")
    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_TEAL)
    page_number(s)


# ── SLIDE 4 ── Macro Problem ──────────────────────────────────────────────────
def slide_04():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_RED)

    text_box(s, Inches(0.8), Inches(1.5), Inches(2), Inches(0.5),
             "PROBLEM", font_size=12, color=ACCENT_RED, bold=True)
    rect(s, Inches(0.8), Inches(2.1), Inches(0.8), Pt(3), fill_color=ACCENT_RED)

    text_box(s, Inches(0.8), Inches(2.6), Inches(11), Inches(0.7),
             "The digital world is breaking trust at scale",
             font_size=32, color=WHITE, bold=True)

    problems = [
        ("\U0001F4F8", "Deepfakes indistinguishable\nfrom real footage"),
        ("\U0001F9D1\u200D\U0001F3A4", "AI-generated personas\nimpersonating individuals"),
        ("\U0001F50D", "Identity theft\nbecoming automated"),
        ("\u2753", "No universal standard for\nverifying \"real human content\""),
    ]
    for i, (icon, text) in enumerate(problems):
        x = Inches(0.8 + i * 3.1)
        y = Inches(3.8)
        rect(s, x, y, Inches(2.8), Inches(2.0), fill_color=RGBColor(0x1E, 0x14, 0x1E), radius=0.05)
        rect(s, x, y, Inches(2.8), Pt(3), fill_color=ACCENT_RED)
        text_box(s, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5),
                 icon, font_size=28, color=ACCENT_RED, font_name="Segoe UI Symbol")
        text_box(s, x + Inches(0.2), y + Inches(0.9), Inches(2.4), Inches(0.9),
                 text, font_size=15, color=LIGHT_GREY, font_name="Calibri")

    # Bottom highlight box
    rect(s, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.8), fill_color=RGBColor(0x1A, 0x0E, 0x0E))
    text_box(s, Inches(1.0), Inches(6.2), Inches(11), Inches(0.6),
             "Today: The internet cannot answer \u201cIs this real or fake?\u201d   |   MYTAR aims to solve exactly this gap.",
             font_size=16, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)
    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_RED)
    page_number(s)


# ── SLIDE 5 ── System Overview ── (Architecture Diagram) ──────────────────────
def slide_05():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_TEAL)
    add_section_header(s, 5, "MYTAR System Overview", "Platform Architecture", accent=False)

    # Remove header elements from section_header since we override
    # Actually, let me just use the raw approach
    # (above we added bg and top bar, so let's redesign this slide cleanly)
    s2 = add_blank_slide()
    bg(s2, BG_DARK)
    rect(s2, 0, 0, W, Pt(3), fill_color=ACCENT_TEAL)
    text_box(s2, Inches(0.8), Inches(0.2), Inches(2), Inches(0.5),
             "05", font_size=14, color=ACCENT_TEAL, bold=True)
    text_box(s2, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "MYTAR System Overview", font_size=32, color=WHITE, bold=True)
    text_box(s2, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "A 4-layer ecosystem built on cryptographic trust",
             font_size=14, color=LIGHT_GREY)
    rect(s2, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_TEAL)

    # ── ARCHITECTURE DIAGRAM ──
    # Vertical stack of 4 layers with connecting arrows + side labels
    layers = [
        ("APPLICATION LAYER", "How you interact with the world", ACCENT_TEAL,
         ["AI Fashion & Commerce", "Personal AI Shopper", "Communication Layer",
          "AR/VR Interfaces", "Meeting Proxy", "Creator Marketplace"]),
        ("AI TWIN LAYER", "How you exist digitally", ACCENT_BLUE,
         ["Voice & Face Models", "Memory Graph", "Behavioral Engine",
          "Knowledge Base", "Context Reasoning", "Multi-modal Output"]),
        ("OWNERSHIP LAYER", "What you own", ACCENT_PURPLE,
         ["NFT Identity Certificate", "Digital Rights Contract",
          "Licensing Engine", "Royalty Distribution", "Usage Audit"]),
        ("IDENTITY LAYER \u2014 Blockchain Core", "Who you are", ACCENT_GREEN,
         ["DID Standards", "Identity Registry", "Immutable Proofs",
          "Crypto Wallet", "Content Signatures"]),
    ]

    left_box = Inches(0.8)
    box_w = Inches(5.5)
    box_h = Inches(1.15)
    gap = Inches(0.15)
    start_y = Inches(2.3)

    # side panel for layer details
    detail_x = Inches(7.0)
    detail_w = Inches(5.8)

    for i, (title, subtitle, color, items) in enumerate(layers):
        y = start_y + i * (box_h + gap)

        # Main layer box
        shape = rect(s2, left_box, y, box_w, box_h, fill_color=RGBColor(0x12, 0x1C, 0x38), radius=0.04)
        # left color accent
        rect(s2, left_box, y, Pt(5), box_h, fill_color=color)
        # layer number
        circle(s2, left_box + Inches(0.25), y + Inches(0.15), Inches(0.35),
               fill_color=color)
        text_box(s2, left_box + Inches(0.25), y + Inches(0.15), Inches(0.35), Inches(0.35),
                 str(4 - i), font_size=14, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        # Layer title
        text_box(s2, left_box + Inches(0.75), y + Inches(0.1), Inches(4.5), Inches(0.35),
                 title, font_size=14, color=color, bold=True)
        # Layer subtitle
        text_box(s2, left_box + Inches(0.75), y + Inches(0.5), Inches(4.5), Inches(0.3),
                 subtitle, font_size=11, color=LIGHT_GREY)

        # Arrow between layers (except last)
        if i < 3:
            arrow_y = y + box_h
            arrow = s2.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left_box + Inches(0.5), arrow_y,
                                        Inches(0.3), gap)
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = color
            arrow.line.fill.background()

        # Detail panel - items grid
        item_x = detail_x
        item_y = y + Inches(0.05)
        items_per_row = 2
        for j, item in enumerate(items):
            ix = item_x + (j % items_per_row) * (detail_w / items_per_row)
            iy = item_y + (j // items_per_row) * Inches(0.32)
            text_box(s2, ix, iy, detail_w / items_per_row - Inches(0.1), Inches(0.3),
                     f"\u25B8 {item}", font_size=10, color=LIGHT_GREY)

    # bottom text
    text_box(s2, left_box, Inches(6.9), Inches(8), Inches(0.4),
             "Each layer builds on cryptographic trust . Identity  \u2192  Ownership  \u2192  AI Twin  \u2192  Applications",
             font_size=12, color=DIM_GREY)

    rect(s2, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_TEAL)
    page_number(s2, 5)
    return s2


# ── SLIDE 6 ── Identity Layer ─────────────────────────────────────────────────
def slide_06():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_GREEN)

    text_box(s, Inches(0.8), Inches(0.2), Inches(2), Inches(0.5),
             "06", font_size=14, color=ACCENT_GREEN, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Identity Layer", font_size=32, color=WHITE, bold=True)
    text_box(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Blockchain backbone for decentralized identity",
             font_size=14, color=LIGHT_GREY)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_GREEN)

    # Left side - key features
    features = [
        ("\U0001F510", "Decentralized Identity (DID)", "W3C-compliant decentralized identifiers"),
        ("\U0001F3D7\ufe0f", "Blockchain Identity Registry", "Immutable on-chain identity anchoring"),
        ("\u2705", "Immutable Identity Proofs", "Tamper-proof verification credentials"),
        ("\U0001F9F0", "Crypto Identity Wallet", "User-owned private key management"),
    ]
    for i, (icon, title, desc) in enumerate(features):
        y = Inches(2.5 + i * 1.1)
        rect(s, Inches(0.8), y, Inches(5.5), Inches(0.9), fill_color=BG_MID, radius=0.04)
        text_box(s, Inches(1.0), y + Inches(0.1), Inches(0.5), Inches(0.5),
                 icon, font_size=22, color=ACCENT_GREEN, font_name="Segoe UI Symbol")
        text_box(s, Inches(1.6), y + Inches(0.1), Inches(4.5), Inches(0.35),
                 title, font_size=15, color=WHITE, bold=True)
        text_box(s, Inches(1.6), y + Inches(0.5), Inches(4.5), Inches(0.3),
                 desc, font_size=11, color=LIGHT_GREY)

    # Right side - hero box
    rect(s, Inches(7.0), Inches(2.5), Inches(5.5), Inches(4.2), fill_color=RGBColor(0x0A, 0x18, 0x14), radius=0.06)
    rect(s, Inches(7.0), Inches(2.5), Inches(5.5), Pt(3), fill_color=ACCENT_GREEN)

    text_box(s, Inches(7.3), Inches(2.8), Inches(5), Inches(0.4),
             "Each identity is:", font_size=14, color=LIGHT_GREY, bold=True)
    text_box(s, Inches(7.3), Inches(3.3), Inches(5), Inches(0.9),
             "A unique cryptographic\nobject anchored on blockchain",
             font_size=22, color=WHITE, bold=True)

    rect(s, Inches(7.3), Inches(4.4), Inches(0.6), Pt(3), fill_color=ACCENT_GREEN)

    props = [
        "\u2713  Tamper-proof",
        "\u2713  Portable across platforms",
        "\u2713  Globally verifiable",
        "\u2713  User-controlled",
    ]
    for j, p in enumerate(props):
        text_box(s, Inches(7.3), Inches(4.7 + j * 0.4), Inches(5), Inches(0.35),
                 p, font_size=14, color=ACCENT_GREEN)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_GREEN)
    page_number(s)


# ── SLIDE 7 ── Identity Verification System ──────────────────────────────────
def slide_07():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_BLUE)

    text_box(s, Inches(0.8), Inches(0.2), Inches(2), Inches(0.5),
             "07", font_size=14, color=ACCENT_BLUE, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Identity Verification System", font_size=32, color=WHITE, bold=True)
    text_box(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Structured onboarding with multi-step verification",
             font_size=14, color=LIGHT_GREY)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_BLUE)

    # Flow: 4 step boxes connected by arrows
    steps = [
        ("\U0001F4CB", "Government ID\nVerification", "Passport, national ID,\ndriver's license"),
        ("\U0001FAE1", "Liveness\nDetection", "Anti-deepfake\nonboarding check"),
        ("\U0001F9EC", "Biometric\nVerification", "User-controlled,\nnot mandatory"),
        ("\u270D\ufe0f", "Multi-step\nConsent Binding", "Cryptographic\nconsent signature"),
    ]
    for i, (icon, title, desc) in enumerate(steps):
        x = Inches(0.8 + i * 3.1)
        y = Inches(2.5)
        rect(s, x, y, Inches(2.8), Inches(2.8), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Inches(2.8), Pt(4), fill_color=ACCENT_BLUE)
        # step number
        circle(s, x + Inches(1.15), y + Inches(0.25), Inches(0.5), fill_color=ACCENT_BLUE)
        text_box(s, x + Inches(1.15), y + Inches(0.25), Inches(0.5), Inches(0.5),
                 str(i + 1), font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # icon
        text_box(s, x + Inches(1.0), y + Inches(0.9), Inches(0.8), Inches(0.5),
                 icon, font_size=28, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER,
                 font_name="Segoe UI Symbol")
        # title
        text_box(s, x + Inches(0.2), y + Inches(1.5), Inches(2.4), Inches(0.6),
                 title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # desc
        text_box(s, x + Inches(0.2), y + Inches(2.1), Inches(2.4), Inches(0.5),
                 desc, font_size=11, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

        # arrow between steps
        if i < 3:
            ax = x + Inches(2.8) + Inches(0.1)
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, y + Inches(1.2),
                                        Inches(0.2), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_BLUE
            arrow.line.fill.background()

    # result box
    rect(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.0), fill_color=RGBColor(0x0E, 0x1A, 0x28), radius=0.04)
    rect(s, Inches(0.8), Inches(5.8), Pt(5), Inches(1.0), fill_color=ACCENT_TEAL)
    text_box(s, Inches(1.5), Inches(5.9), Inches(10), Inches(0.8),
             "Once verified: The identity is anchored on-chain as a permanent digital certificate.",
             font_size=16, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_BLUE)
    page_number(s)


# ── SLIDE 8 ── NFT-Based Identity Ownership ──────────────────────────────────
def slide_08():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_PURPLE)

    text_box(s, Inches(0.8), Inches(0.2), Inches(2), Inches(0.5),
             "08", font_size=14, color=ACCENT_PURPLE, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "NFT-Based Identity Ownership", font_size=32, color=WHITE, bold=True)
    text_box(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Functional NFT as a digital identity deed",
             font_size=14, color=LIGHT_GREY)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_PURPLE)

    # Centerpiece: large NFT visual representation
    nft_box = rect(s, Inches(1.0), Inches(2.5), Inches(4.5), Inches(4.2),
                   fill_color=RGBColor(0x1C, 0x14, 0x28), radius=0.06)
    rect(s, Inches(1.0), Inches(2.5), Inches(4.5), Pt(4), fill_color=ACCENT_PURPLE)

    text_box(s, Inches(1.3), Inches(2.8), Inches(4), Inches(0.4),
             "IDENTITY NFT", font_size=12, color=ACCENT_PURPLE, bold=True)

    # Diamond icon
    diamond = s.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(2.4), Inches(3.4), Inches(1.2), Inches(1.0))
    diamond.fill.solid()
    diamond.fill.fore_color.rgb = ACCENT_PURPLE
    diamond.line.fill.background()
    text_box(s, Inches(2.4), Inches(3.5), Inches(1.2), Inches(0.8),
             "\U0001F511", font_size=36, color=WHITE, alignment=PP_ALIGN.CENTER,
             font_name="Segoe UI Symbol")

    text_box(s, Inches(1.3), Inches(4.7), Inches(4), Inches(0.3),
             "Ownership of digital identity", font_size=14, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
    text_box(s, Inches(1.3), Inches(5.1), Inches(4), Inches(0.3),
             "Rights over digital likeness", font_size=13, color=LIGHT_GREY,
             alignment=PP_ALIGN.CENTER)
    text_box(s, Inches(1.3), Inches(5.5), Inches(4), Inches(0.3),
             "Permissioned usage of avatar & content", font_size=13, color=LIGHT_GREY,
             alignment=PP_ALIGN.CENTER)
    text_box(s, Inches(1.3), Inches(5.9), Inches(4), Inches(0.3),
             "Licensing capabilities", font_size=13, color=LIGHT_GREY,
             alignment=PP_ALIGN.CENTER)

    # Right side - explanation
    rect(s, Inches(6.0), Inches(2.5), Inches(6.5), Inches(4.2), fill_color=BG_MID, radius=0.04)

    text_box(s, Inches(6.3), Inches(2.8), Inches(6), Inches(0.4),
             "This NFT is not speculative \u2014 it is functional.", font_size=16,
             color=ACCENT_TEAL, bold=True)
    rect(s, Inches(6.3), Inches(3.3), Inches(0.6), Pt(3), fill_color=ACCENT_TEAL)

    roles = [
        "Identity Deed \u2014 Proof of personhood",
        "Digital Rights Contract \u2014 Usage control",
        "Licensing Instrument \u2014 Monetization",
        "Authentication Token \u2014 Verification",
    ]
    for j, role in enumerate(roles):
        text_box(s, Inches(6.3), Inches(3.6 + j * 0.6), Inches(6), Inches(0.5),
                 f"\u25B8  {role}", font_size=14, color=LIGHT_GREY)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_PURPLE)
    page_number(s)


# ── SLIDE 9 ── Blockchain Utility Layer ──────────────────────────────────────
def slide_09():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_TEAL)

    text_box(s, Inches(0.8), Inches(0.2), Inches(2), Inches(0.5),
             "09", font_size=14, color=ACCENT_TEAL, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Blockchain Utility Layer", font_size=32, color=WHITE, bold=True)
    text_box(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Transparent, tamper-proof digital human rights infrastructure",
             font_size=14, color=LIGHT_GREY)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_TEAL)

    utilities = [
        ("\U0001F4C1", "Immutable Identity Records", "Permanent, unchangeable\nidentity anchoring"),
        ("\U0001F50D", "Content Provenance Tracking", "Trace origin of every\ndigital artifact"),
        ("\U0001F4B0", "Licensing & Royalties", "Automated distribution\nvia smart contracts"),
        ("\U0001F4CA", "Avatar Usage Auditing", "Track every use of\nyour digital likeness"),
        ("\u2705", "Consent Verification Logs", "Cryptographic proof\nof permission grants"),
    ]
    for i, (icon, title, desc) in enumerate(utilities):
        x = Inches(0.8 + (i % 3) * 4.1)
        y = Inches(2.5 + (i // 3) * 2.1)
        rect(s, x, y, Inches(3.7), Inches(1.8), fill_color=BG_MID, radius=0.04)
        rect(s, x, y, Pt(4), Inches(1.8), fill_color=ACCENT_TEAL)
        text_box(s, x + Inches(0.3), y + Inches(0.2), Inches(0.5), Inches(0.5),
                 icon, font_size=24, color=ACCENT_TEAL, font_name="Segoe UI Symbol")
        text_box(s, x + Inches(0.9), y + Inches(0.2), Inches(2.5), Inches(0.35),
                 title, font_size=15, color=WHITE, bold=True)
        text_box(s, x + Inches(0.9), y + Inches(0.7), Inches(2.5), Inches(0.8),
                 desc, font_size=11, color=LIGHT_GREY)

    # Bottom statement
    rect(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.6), fill_color=RGBColor(0x0E, 0x1C, 0x18), radius=0.04)
    text_box(s, Inches(1.0), Inches(6.55), Inches(11), Inches(0.5),
             "A transparent, tamper-proof system of digital human rights.",
             font_size=14, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_TEAL)
    page_number(s)


# ── SLIDE 10 ── Content Authenticity Layer ───────────────────────────────────
def slide_10():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_BLUE)

    text_box(s, Inches(0.8), Inches(0.2), Inches(2), Inches(0.5),
             "10", font_size=14, color=ACCENT_BLUE, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Content Authenticity Layer", font_size=32, color=WHITE, bold=True)
    text_box(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Cryptographic signing of digital artifacts",
             font_size=14, color=LIGHT_GREY)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_BLUE)

    # 3 column layout
    cols = [
        ("\U0001F4DD", "Signing", "Every digital artifact\n(video, image, voice,\nAI output) can be:\n\n\u25B8 Signed using identity keys\n\u25B8 Timestamped on blockchain\n\u25B8 Verified as authentic\n  or synthetic"),
        ("\u2696\ufe0f", "Verification", "The system enables:\n\n\u25B8 Proof of real human content\n\u25B8 Detection of impersonation\n\u25B8 Traceability of media origin\n\u25B8 Trust scoring of content"),
        ("\U0001F3ED", "Ecosystem", "Built for:\n\n\u25B8 Social media platforms\n\u25B8 News & journalism\n\u25B8 Legal evidence\n\u25B8 Content creation\n\u25B8 Enterprise communications"),
    ]
    for i, (icon, title, desc) in enumerate(cols):
        x = Inches(0.8 + i * 4.1)
        y = Inches(2.5)
        rect(s, x, y, Inches(3.7), Inches(4.2), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Inches(3.7), Pt(4), fill_color=ACCENT_BLUE)
        text_box(s, x + Inches(0.3), y + Inches(0.3), Inches(0.5), Inches(0.5),
                 icon, font_size=28, color=ACCENT_BLUE, font_name="Segoe UI Symbol")
        text_box(s, x + Inches(0.9), y + Inches(0.3), Inches(2.5), Inches(0.4),
                 title, font_size=18, color=WHITE, bold=True)
        text_box(s, x + Inches(0.3), y + Inches(1.1), Inches(3.1), Inches(2.8),
                 desc, font_size=12, color=LIGHT_GREY)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_BLUE)
    page_number(s)


# ── SLIDE 11 ── Digital Twin Engine ──────────────────────────────────────────
def slide_11():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_TEAL)

    text_box(s, Inches(0.8), Inches(0.2), Inches(2), Inches(0.5),
             "11", font_size=14, color=ACCENT_TEAL, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Digital Twin Engine", font_size=32, color=WHITE, bold=True)
    text_box(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "A continuously evolving AI representation of your identity",
             font_size=14, color=LIGHT_GREY)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_TEAL)

    # 2x2 grid of twin components
    components = [
        ("\U0001F399\ufe0f", "Voice Replica Model", "Your vocal fingerprint\nsynthesized by AI"),
        ("\U0001F5E3\ufe0f", "Facial Representation", "Photorealistic digital\nface model"),
        ("\U0001F9E0", "Behavioral Pattern Model", "Decision patterns &\ncommunication style"),
        ("\U0001F4BE", "Memory Graph", "Facts, preferences,\nhistory & knowledge"),
    ]
    for i, (icon, title, desc) in enumerate(components):
        x = Inches(0.8 + (i % 2) * 6.0)
        y = Inches(2.5 + (i // 2) * 2.0)
        rect(s, x, y, Inches(5.5), Inches(1.7), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Pt(5), Inches(1.7), fill_color=ACCENT_TEAL)
        text_box(s, x + Inches(0.3), y + Inches(0.3), Inches(0.6), Inches(0.6),
                 icon, font_size=28, color=ACCENT_TEAL, font_name="Segoe UI Symbol")
        text_box(s, x + Inches(1.1), y + Inches(0.2), Inches(4), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True)
        text_box(s, x + Inches(1.1), y + Inches(0.7), Inches(4), Inches(0.8),
                 desc, font_size=13, color=LIGHT_GREY)

    # Bottom emphasis
    rect(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8), fill_color=RGBColor(0x0E, 0x1C, 0x18), radius=0.04)
    text_box(s, Inches(1.0), Inches(6.3), Inches(11), Inches(0.6),
             "This is not just an avatar. It is a computational representation of your identity.",
             font_size=16, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_TEAL)
    page_number(s)


# ── SLIDE 12 ── AI Brain Layer ───────────────────────────────────────────────
def slide_12():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_BLUE)

    text_box(s, Inches(0.8), Inches(0.2), Inches(2), Inches(0.5),
             "12", font_size=14, color=ACCENT_BLUE, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "AI Brain Layer", font_size=32, color=WHITE, bold=True)
    text_box(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Intelligence engine powering the digital twin",
             font_size=14, color=LIGHT_GREY)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_BLUE)

    # Brain architecture - central box + 4 surrounding modules
    # Central
    rect(s, Inches(4.2), Inches(3.2), Inches(5), Inches(1.5), fill_color=RGBColor(0x14, 0x20, 0x30), radius=0.06)
    rect(s, Inches(4.2), Inches(3.2), Inches(5), Pt(4), fill_color=ACCENT_TEAL)
    text_box(s, Inches(4.5), Inches(3.4), Inches(4.5), Inches(0.4),
             "AI TWIN CORE", font_size=16, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)
    text_box(s, Inches(4.5), Inches(3.8), Inches(4.5), Inches(0.6),
             "Responds as \"you would respond\"",
             font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Surrounding modules
    modules = [
        (Inches(0.8), Inches(2.5), Inches(3.0), Inches(1.2), ACCENT_BLUE,
         "Personal Knowledge Base", "Documents, chats,\ninteractions"),
        (Inches(9.5), Inches(2.5), Inches(3.0), Inches(1.2), ACCENT_PURPLE,
         "Context-Aware\nReasoning Engine", "Situational\nunderstanding"),
        (Inches(0.8), Inches(5.0), Inches(3.0), Inches(1.2), ACCENT_GREEN,
         "Long-Term Memory\nSystem", "Persistent\nknowledge retention"),
        (Inches(9.5), Inches(5.0), Inches(3.0), Inches(1.2), ACCENT_ORANGE,
         "Behavioral Adaptation\nModels", "Style & preference\nlearning"),
    ]
    for (mx, my, mw, mh, mcolor, mtitle, mdesc) in modules:
        rect(s, mx, my, mw, mh, fill_color=BG_MID, radius=0.05)
        rect(s, mx, my, mw, Pt(3), fill_color=mcolor)
        text_box(s, mx + Inches(0.15), my + Inches(0.15), mw - Inches(0.3), Inches(0.45),
                 mtitle, font_size=12, color=WHITE, bold=True)
        text_box(s, mx + Inches(0.15), my + Inches(0.6), mw - Inches(0.3), Inches(0.5),
                 mdesc, font_size=10, color=LIGHT_GREY)

    # Arrows from modules to center (simple lines via thin rects)
    # (skipping arrows for simplicity, using visual position)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_BLUE)
    page_number(s)


# ── SLIDE 13 ── Interaction Modes ────────────────────────────────────────────
def slide_13():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_PURPLE)

    text_box(s, Inches(0.8), Inches(0.2), Inches(2), Inches(0.5),
             "13", font_size=14, color=ACCENT_PURPLE, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Interaction Modes", font_size=32, color=WHITE, bold=True)
    text_box(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Multi-interface AI twin presence",
             font_size=14, color=LIGHT_GREY)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_PURPLE)

    modes = [
        ("\U0001F4F9", "Video Avatar", "Face + voice simulation\nfor real-time interaction",
         ACCENT_TEAL),
        ("\U0001F399\ufe0f", "Voice Assistant", "Hands-free interaction\nacross devices",
         ACCENT_BLUE),
        ("\U0001F4AC", "Chat Assistant", "Text-based twin\nfor messaging",
         ACCENT_GREEN),
        ("\U0001F91D", "Meeting Proxy", "Attends meetings\non your behalf",
         ACCENT_ORANGE),
    ]
    for i, (icon, title, desc, color) in enumerate(modes):
        x = Inches(0.8 + i * 3.1)
        y = Inches(2.5)
        rect(s, x, y, Inches(2.8), Inches(3.0), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Inches(2.8), Pt(5), fill_color=color)
        text_box(s, x + Inches(0.2), y + Inches(0.3), Inches(2.4), Inches(0.6),
                 icon, font_size=36, color=color, alignment=PP_ALIGN.CENTER,
                 font_name="Segoe UI Symbol")
        text_box(s, x + Inches(0.2), y + Inches(1.1), Inches(2.4), Inches(0.5),
                 title, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        text_box(s, x + Inches(0.2), y + Inches(1.7), Inches(2.4), Inches(0.8),
                 desc, font_size=12, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

    # Bottom text
    text_box(s, Inches(0.8), Inches(5.9), Inches(11), Inches(0.5),
             "Each mode is governed by identity permissions and cryptographic consent.",
             font_size=14, color=DIM_GREY, alignment=PP_ALIGN.CENTER)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_PURPLE)
    page_number(s)


# ── SLIDE 14 ── AI Fashion & Identity Commerce ───────────────────────────────
def slide_14():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_TEAL)

    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "14  |  USE CASE", font_size=14, color=ACCENT_TEAL, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "AI Fashion & Identity Commerce", font_size=32, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_TEAL)

    # Two column layout
    # Left: feature list
    features = [
        ("\U0001F457", "Persistent Body Model", "User body dimensions stored securely"),
        ("\U0001F484", "AI Fashion Recommendations", "Personalized style suggestions"),
        ("\U0001F4F1", "Virtual Try-On Systems", "AR-powered fitting experience"),
        ("\U0001F3F7", "Brand Integrations", "Zara, H&M, Nike, and more"),
    ]
    for i, (icon, title, desc) in enumerate(features):
        y = Inches(2.5 + i * 1.1)
        rect(s, Inches(0.8), y, Inches(5.5), Inches(0.9), fill_color=BG_MID, radius=0.04)
        text_box(s, Inches(1.0), y + Inches(0.15), Inches(0.5), Inches(0.5),
                 icon, font_size=22, color=ACCENT_TEAL, font_name="Segoe UI Symbol")
        text_box(s, Inches(1.6), y + Inches(0.1), Inches(4.5), Inches(0.35),
                 title, font_size=15, color=WHITE, bold=True)
        text_box(s, Inches(1.6), y + Inches(0.5), Inches(4.5), Inches(0.3),
                 desc, font_size=11, color=LIGHT_GREY)

    # Right: differentiation box
    rect(s, Inches(7.0), Inches(2.5), Inches(5.5), Inches(4.2),
         fill_color=RGBColor(0x12, 0x1C, 0x28), radius=0.06)
    rect(s, Inches(7.0), Inches(2.5), Inches(5.5), Pt(3), fill_color=ACCENT_TEAL)
    text_box(s, Inches(7.3), Inches(2.8), Inches(5), Inches(0.4),
             "Unlike traditional systems:", font_size=14, color=ACCENT_TEAL, bold=True)
    text_box(s, Inches(7.3), Inches(3.3), Inches(5), Inches(1.5),
             "The model persists\nacross platforms.\n\nYour digital body is yours,\nnot locked into any\ne-commerce silo.",
             font_size=18, color=WHITE, bold=True)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_TEAL)
    page_number(s)


# ── SLIDE 15 ── Personal AI Shopper ──────────────────────────────────────────
def slide_15():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_BLUE)

    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "15  |  USE CASE", font_size=14, color=ACCENT_BLUE, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Personal AI Shopper", font_size=32, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_BLUE)

    features = [
        ("\U0001F52E", "Predicts Purchase Intent", "Anticipates needs\nbefore you search"),
        ("\U0001F455", "Suggests Outfits & Products", "Curated recommendations\nbased on style profile"),
        ("\U0001F9D1\u200D\U0001F4BB", "Learns Preferences", "Continuous learning\nfrom your choices"),
        ("\U0001F310", "E-Commerce Integration", "Seamless API connections\nto major retailers"),
    ]
    for i, (icon, title, desc) in enumerate(features):
        x = Inches(0.8 + i * 3.1)
        y = Inches(2.5)
        rect(s, x, y, Inches(2.8), Inches(2.5), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Inches(2.8), Pt(4), fill_color=ACCENT_BLUE)
        text_box(s, x + Inches(0.2), y + Inches(0.3), Inches(2.4), Inches(0.5),
                 icon, font_size=28, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER,
                 font_name="Segoe UI Symbol")
        text_box(s, x + Inches(0.2), y + Inches(1.0), Inches(2.4), Inches(0.4),
                 title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        text_box(s, x + Inches(0.2), y + Inches(1.5), Inches(2.4), Inches(0.8),
                 desc, font_size=11, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

    # Bottom tag
    rect(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.7), fill_color=RGBColor(0x0E, 0x18, 0x28), radius=0.04)
    text_box(s, Inches(1.0), Inches(5.6), Inches(11), Inches(0.5),
             "This creates: Identity-driven commerce personalization at scale.",
             font_size=16, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_BLUE)
    page_number(s)


# ── SLIDE 16 ── Body & Appearance Simulation ────────────────────────────────
def slide_16():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_PURPLE)
    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "16  |  USE CASE", font_size=14, color=ACCENT_PURPLE, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Body & Appearance Simulation", font_size=32, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_PURPLE)

    use_cases = [
        ("\U0001FA78", "Cosmetic Surgery\nPreviews", "Visualize outcomes\nbefore procedures"),
        ("\U0001F9CD", "Facial\nModifications", "Simulate aesthetic\nchanges in real-time"),
        ("\U0001F4AA", "Body Transformation\nModeling", "See fitness &\nphysique changes"),
        ("\U0001FA7A", "Dermatological\nOutcomes", "Skin treatment\npreview system"),
    ]
    for i, (icon, title, desc) in enumerate(use_cases):
        x = Inches(0.8 + i * 3.1)
        y = Inches(2.5)
        rect(s, x, y, Inches(2.8), Inches(2.8), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Inches(2.8), Pt(4), fill_color=ACCENT_PURPLE)
        text_box(s, x + Inches(0.2), y + Inches(0.3), Inches(2.4), Inches(0.6),
                 icon, font_size=32, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER,
                 font_name="Segoe UI Symbol")
        text_box(s, x + Inches(0.2), y + Inches(1.1), Inches(2.4), Inches(0.6),
                 title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        text_box(s, x + Inches(0.2), y + Inches(1.8), Inches(2.4), Inches(0.8),
                 desc, font_size=11, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

    text_box(s, Inches(0.8), Inches(5.7), Inches(11), Inches(0.5),
             "This connects AI with real-world medical decision support.",
             font_size=14, color=DIM_GREY, alignment=PP_ALIGN.CENTER)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_PURPLE)
    page_number(s)


# ── SLIDE 17 ── Digital Proxy (UTAH System) ──────────────────────────────────
def slide_17():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_TEAL)

    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "17  |  USE CASE", font_size=14, color=ACCENT_TEAL, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Digital Proxy  (UTAH System)", font_size=32, color=WHITE, bold=True)
    text_box(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Delegate your digital presence to your AI twin",
             font_size=14, color=LIGHT_GREY)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_TEAL)

    capabilities = [
        ("\U0001F91D", "Attends meetings as\nyour AI twin", "Calendar integration, note-taking,\nfollow-up actions"),
        ("\U0001F4E8", "Responds to emails\nand messages", "Style-matched replies,\nconsent-verified"),
        ("\U0001F4C5", "Handles scheduling\nand communication", "Autonomous calendar\nmanagement"),
        ("\u2699\ufe0f", "Executes predefined\ndecision rules", "Rule-based autonomous\noperation"),
    ]
    for i, (icon, title, desc) in enumerate(capabilities):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(2.3 + (i // 2) * 2.0)
        rect(s, x, y, Inches(5.7), Inches(1.7), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Pt(5), Inches(1.7), fill_color=ACCENT_TEAL)
        text_box(s, x + Inches(0.3), y + Inches(0.3), Inches(0.6), Inches(0.6),
                 icon, font_size=28, color=ACCENT_TEAL, font_name="Segoe UI Symbol")
        text_box(s, x + Inches(1.1), y + Inches(0.2), Inches(4.3), Inches(0.45),
                 title, font_size=16, color=WHITE, bold=True)
        text_box(s, x + Inches(1.1), y + Inches(0.8), Inches(4.3), Inches(0.7),
                 desc, font_size=12, color=LIGHT_GREY)

    rect(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.7), fill_color=RGBColor(0x0E, 0x1C, 0x18), radius=0.04)
    text_box(s, Inches(1.0), Inches(6.3), Inches(11), Inches(0.5),
             "This creates: A scalable extension of human productivity.",
             font_size=16, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_TEAL)
    page_number(s)


# ── SLIDE 18 ── Global Communication Layer ───────────────────────────────────
def slide_18():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_BLUE)

    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "18  |  USE CASE", font_size=14, color=ACCENT_BLUE, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Global Communication Layer", font_size=32, color=WHITE, bold=True)
    text_box(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Breaking language and cultural barriers with AI",
             font_size=14, color=LIGHT_GREY)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_BLUE)

    features = [
        ("\U0001F30D", "Real-Time Translation", "Seamless language conversion\nvia AI twin"),
        ("\U0001F30F", "Cultural Adaptation", "Context-aware communication\nstyle adjustment"),
        ("\U0001F310", "Cross-Border Identity", "Unified identity across\ngeographic boundaries"),
        ("\U0001F4E1", "Universal Interface", "Single communication layer\nfor global interaction"),
    ]
    for i, (icon, title, desc) in enumerate(features):
        x = Inches(0.8 + i * 3.1)
        y = Inches(2.5)
        rect(s, x, y, Inches(2.8), Inches(2.5), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Inches(2.8), Pt(4), fill_color=ACCENT_BLUE)
        text_box(s, x + Inches(0.2), y + Inches(0.3), Inches(2.4), Inches(0.5),
                 icon, font_size=28, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER,
                 font_name="Segoe UI Symbol")
        text_box(s, x + Inches(0.2), y + Inches(1.0), Inches(2.4), Inches(0.4),
                 title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        text_box(s, x + Inches(0.2), y + Inches(1.5), Inches(2.4), Inches(0.8),
                 desc, font_size=11, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_BLUE)
    page_number(s)


# ── SLIDE 19 ── Creator Economy & Avatar Licensing ──────────────────────────
def slide_19():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_PURPLE)

    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "19  |  USE CASE", font_size=14, color=ACCENT_PURPLE, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Creator Economy & Avatar Licensing", font_size=32, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_PURPLE)

    # Left: monetization methods
    methods = [
        ("\U0001F3AC", "License Avatar for Content", "Allow your twin to star\nin digital productions"),
        ("\U0001F399\ufe0f", "Rent Voice & Likeness", "Time-bound licensing\nof your digital self"),
        ("\U0001F31F", "AI-Generated Endorsements", "Brand partnerships via\nyour authorized twin"),
        ("\U0001F4B0", "Blockchain-Tracked Royalties", "Automatic payment\nfor every use"),
    ]
    for i, (icon, title, desc) in enumerate(methods):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(2.3 + (i // 2) * 1.5)
        rect(s, x, y, Inches(5.7), Inches(1.3), fill_color=BG_MID, radius=0.04)
        text_box(s, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5),
                 icon, font_size=22, color=ACCENT_PURPLE, font_name="Segoe UI Symbol")
        text_box(s, x + Inches(0.9), y + Inches(0.15), Inches(2.5), Inches(0.35),
                 title, font_size=15, color=WHITE, bold=True)
        text_box(s, x + Inches(0.9), y + Inches(0.55), Inches(4.5), Inches(0.6),
                 desc, font_size=11, color=LIGHT_GREY)

    # Right: principle box
    rect(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.2), fill_color=RGBColor(0x1C, 0x14, 0x28), radius=0.04)
    rect(s, Inches(0.8), Inches(5.5), Pt(5), Inches(1.2), fill_color=ACCENT_PURPLE)
    text_box(s, Inches(1.5), Inches(5.6), Inches(10), Inches(1.0),
             "All usage is:  Transparent  \u2022  Permissioned via NFT identity contracts  \u2022  Tracked on blockchain",
             font_size=16, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_PURPLE)
    page_number(s)


# ── SLIDE 20 ── Digital Legacy System ────────────────────────────────────────
def slide_20():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_TEAL)

    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "20  |  USE CASE", font_size=14, color=ACCENT_TEAL, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Digital Legacy System", font_size=32, color=WHITE, bold=True)
    text_box(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Preserving human identity beyond lifespan",
             font_size=14, color=LIGHT_GREY)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_TEAL)

    legacy_items = [
        ("\U0001F9E0", "AI Memory\nReconstruction", "Preserve memories and\nknowledge structures"),
        ("\U0001F5E3\ufe0f", "Conversational\nLegacy Avatars", "Interactive digital\npresence for loved ones"),
        ("\U0001F510", "Controlled Access\nby Family/Trustees", "Permissioned access\nto digital inheritance"),
        ("\U0001F4DA", "Knowledge\nPreservation System", "Store wisdom, stories,\nand life experience"),
    ]
    for i, (icon, title, desc) in enumerate(legacy_items):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(2.3 + (i // 2) * 2.0)
        rect(s, x, y, Inches(5.7), Inches(1.7), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Pt(5), Inches(1.7), fill_color=ACCENT_TEAL)
        text_box(s, x + Inches(0.3), y + Inches(0.3), Inches(0.6), Inches(0.6),
                 icon, font_size=28, color=ACCENT_TEAL, font_name="Segoe UI Symbol")
        text_box(s, x + Inches(1.1), y + Inches(0.2), Inches(4.3), Inches(0.5),
                 title, font_size=16, color=WHITE, bold=True)
        text_box(s, x + Inches(1.1), y + Inches(0.8), Inches(4.3), Inches(0.7),
                 desc, font_size=12, color=LIGHT_GREY)

    rect(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.7), fill_color=RGBColor(0x0E, 0x1C, 0x18), radius=0.04)
    text_box(s, Inches(1.0), Inches(6.3), Inches(11), Inches(0.5),
             "A structured digital continuation of human identity.",
             font_size=16, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_TEAL)
    page_number(s)


# ── SLIDE 21 ── Share Your Tar ───────────────────────────────────────────────
def slide_21():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_BLUE)

    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "21  |  USE CASE", font_size=14, color=ACCENT_BLUE, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Communication Layer  (\u201cShare Your Tar\u201d)", font_size=32, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_BLUE)

    comm_features = [
        ("\U0001F916\u200D\U0001F4F1", "Avatar-to-Avatar\nMessaging", "Your twin talks to\ntheir twin"),
        ("\u2705", "Identity-Verified\nCommunication", "Every message\ncryptographically signed"),
        ("\U0001F3D7\ufe0f", "Blockchain-Authenticated\nInteractions", "Immutable proof of\ncommunication"),
        ("\U0001F465", "Secure Identity-Based\nSocial Graph", "Trusted connections\nonly"),
    ]
    for i, (icon, title, desc) in enumerate(comm_features):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(2.3 + (i // 2) * 2.0)
        rect(s, x, y, Inches(5.7), Inches(1.7), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Pt(5), Inches(1.7), fill_color=ACCENT_BLUE)
        text_box(s, x + Inches(0.3), y + Inches(0.3), Inches(0.6), Inches(0.6),
                 icon, font_size=28, color=ACCENT_BLUE, font_name="Segoe UI Symbol")
        text_box(s, x + Inches(1.1), y + Inches(0.2), Inches(4.3), Inches(0.5),
                 title, font_size=16, color=WHITE, bold=True)
        text_box(s, x + Inches(1.1), y + Inches(0.8), Inches(4.3), Inches(0.7),
                 desc, font_size=12, color=LIGHT_GREY)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_BLUE)
    page_number(s)


# ── SLIDE 22 ── Future Interfaces ────────────────────────────────────────────
def slide_22():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_PURPLE)

    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "22  |  FUTURE", font_size=14, color=ACCENT_PURPLE, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Future Interfaces  (XR / Hologram Layer)", font_size=32, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_PURPLE)

    futures = [
        ("\U0001F4F1", "AR-Based Presence", "Augmented reality\nidentity overlay"),
        ("\U0001F3A0", "VR Immersive Identity", "Full virtual reality\nrepresentation"),
        ("\U0001F4E1", "Holographic\nCommunication", "3D holographic\ninteraction systems"),
        ("\U0001F3AE", "Metaverse-Compatible\nIdentity Objects", "Cross-platform\navatar portability"),
    ]
    for i, (icon, title, desc) in enumerate(futures):
        x = Inches(0.8 + i * 3.1)
        y = Inches(2.5)
        rect(s, x, y, Inches(2.8), Inches(2.8), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Inches(2.8), Pt(4), fill_color=ACCENT_PURPLE)
        text_box(s, x + Inches(0.2), y + Inches(0.3), Inches(2.4), Inches(0.6),
                 icon, font_size=32, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER,
                 font_name="Segoe UI Symbol")
        text_box(s, x + Inches(0.2), y + Inches(1.1), Inches(2.4), Inches(0.6),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        text_box(s, x + Inches(0.2), y + Inches(1.8), Inches(2.4), Inches(0.8),
                 desc, font_size=12, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_PURPLE)
    page_number(s)


# ── SLIDE 23 ── Competitive Landscape ────────────────────────────────────────
def slide_23():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_ORANGE)

    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "23  |  MARKET", font_size=14, color=ACCENT_ORANGE, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Competitive Landscape", font_size=32, color=WHITE, bold=True)
    text_box(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "The market is fragmented across multiple categories",
             font_size=14, color=LIGHT_GREY)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_ORANGE)

    categories = [
        ("Identity Systems", ["World ID (Tools for Humanity)", "Microsoft Entra Verified ID"], ACCENT_TEAL),
        ("Content Authenticity", ["Adobe Content Credentials", "Google SynthID"], ACCENT_BLUE),
        ("AI Avatars", ["HeyGen", "Synthesia", "D-ID"], ACCENT_PURPLE),
        ("AI Companions", ["Replika", "Character AI"], ACCENT_GREEN),
    ]
    for i, (cat_title, players, color) in enumerate(categories):
        x = Inches(0.8 + i * 3.1)
        y = Inches(2.3)
        rect(s, x, y, Inches(2.8), Inches(2.8), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Inches(2.8), Pt(4), fill_color=color)
        text_box(s, x + Inches(0.2), y + Inches(0.2), Inches(2.4), Inches(0.4),
                 cat_title, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        for j, player in enumerate(players):
            text_box(s, x + Inches(0.2), y + Inches(0.8 + j * 0.5), Inches(2.4), Inches(0.4),
                     f"\u25B8  {player}", font_size=11, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

    # MYTAR differentiation
    rect(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.2), fill_color=RGBColor(0x1C, 0x14, 0x0E), radius=0.04)
    rect(s, Inches(0.8), Inches(5.5), Pt(5), Inches(1.2), fill_color=ACCENT_ORANGE)
    text_box(s, Inches(1.5), Inches(5.6), Inches(10.5), Inches(1.0),
             "None of these combine: Identity + Ownership + AI Twin + Blockchain + Commerce  |  MYTAR is the first unified platform.",
             font_size=16, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_ORANGE)
    page_number(s)


# ── SLIDE 24 ── Market Fragmentation ──────────────────────────────────────────
def slide_24():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_ORANGE)
    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "24  |  COMPETITIVE DEEP DIVE", font_size=14, color=ACCENT_ORANGE, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Market Fragmentation", font_size=32, color=WHITE, bold=True)
    text_box(s, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
             "The market is highly fragmented across four emerging but disconnected categories",
             font_size=14, color=LIGHT_GREY)
    rect(s, Inches(0.8), Inches(1.9), Inches(1.5), Pt(3), fill_color=ACCENT_ORANGE)

    # 4 category boxes - showing fragmentation
    cats = [
        ("\U0001F511", "Identity Systems", "\u201cWho are you?\u201d", ACCENT_TEAL),
        ("\u2705", "Content Authenticity", "\u201cIs this real?\u201d", ACCENT_BLUE),
        ("\U0001F5E3\ufe0f", "AI Avatar Platforms", "\u201cWhat do you look like?\u201d", ACCENT_PURPLE),
        ("\U0001F9D1\u200D\U0001F4BB", "AI Companions", "\u201cWho interacts with you?\u201d", ACCENT_GREEN),
    ]
    for i, (icon, title, question, color) in enumerate(cats):
        x = Inches(0.8 + i * 3.1)
        y = Inches(2.5)
        rect(s, x, y, Inches(2.8), Inches(2.5), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Inches(2.8), Pt(4), fill_color=color)
        text_box(s, x + Inches(0.2), y + Inches(0.3), Inches(2.4), Inches(0.5),
                 icon, font_size=28, color=color, alignment=PP_ALIGN.CENTER,
                 font_name="Segoe UI Symbol")
        text_box(s, x + Inches(0.2), y + Inches(1.0), Inches(2.4), Inches(0.4),
                 title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        text_box(s, x + Inches(0.2), y + Inches(1.5), Inches(2.4), Inches(0.8),
                 question, font_size=12, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

    # Bottom context
    rect(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.3), fill_color=RGBColor(0x1C, 0x14, 0x0E), radius=0.04)
    rect(s, Inches(0.8), Inches(5.5), Pt(5), Inches(1.3), fill_color=ACCENT_ORANGE)
    text_box(s, Inches(1.5), Inches(5.6), Inches(10.5), Inches(1.1),
             "Despite rapid innovation in digital identity, AI avatars, and content authenticity,\ntoday\u2019s ecosystem is split into isolated vertical stacks,\neach solving only a narrow part of the problem.",
             font_size=14, color=LIGHT_GREY)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_ORANGE)
    page_number(s)


# ── SLIDE 25 ── Identity Systems Deep Dive ────────────────────────────────────
def slide_25():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_TEAL)
    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "25  |  COMPETITOR ANALYSIS", font_size=14, color=ACCENT_TEAL, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Identity Systems  (\u201cWho are you?\u201d)", font_size=32, color=WHITE, bold=True)
    text_box(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "These platforms focus on verifying human identity or credentials",
             font_size=14, color=LIGHT_GREY)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_TEAL)

    # Two-column competitor cards
    comps = [
        ("World ID", "Tools for Humanity", ACCENT_TEAL, [
            ("Focus", "Proof-of-human via biometric scanning /\niris-based identity"),
            ("Strength", "Strong anti-bot / Sybil resistance\nfor digital platforms"),
            ("Limitation", "No concept of AI persona, content\nownership, or monetization layer"),
        ]),
        ("Microsoft Entra Verified ID", "Enterprise verifiable credentials", ACCENT_BLUE, [
            ("Focus", "Enterprise-grade verifiable credentials\n(work, education, access)"),
            ("Strength", "Deep enterprise integration,\ncompliance-ready identity verification"),
            ("Limitation", "Static identity layer; not designed for\nconsumer AI or creative identity extension"),
        ]),
    ]
    for ci, (name, subtitle, color, details) in enumerate(comps):
        x = Inches(0.8 + ci * 6.2)
        y = Inches(2.3)
        rect(s, x, y, Inches(5.7), Inches(4.0), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Inches(5.7), Pt(4), fill_color=color)
        text_box(s, x + Inches(0.3), y + Inches(0.2), Inches(5.1), Inches(0.35),
                 name, font_size=18, color=WHITE, bold=True)
        text_box(s, x + Inches(0.3), y + Inches(0.6), Inches(5.1), Inches(0.3),
                 subtitle, font_size=11, color=DIM_GREY)
        for j, (label, text) in enumerate(details):
            dy = y + Inches(1.0 + j * 1.0)
            text_box(s, x + Inches(0.3), dy, Inches(1.2), Inches(0.25),
                     label, font_size=10, color=color, bold=True)
            text_box(s, x + Inches(0.3), dy + Inches(0.25), Inches(5.1), Inches(0.65),
                     text, font_size=12, color=LIGHT_GREY)

    # Key gap
    rect(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6), fill_color=RGBColor(0x0E, 0x1C, 0x18), radius=0.04)
    text_box(s, Inches(1.0), Inches(6.65), Inches(11), Inches(0.5),
             "Key gap: Identity is treated as a verification problem, not a living digital asset.",
             font_size=14, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_TEAL)
    page_number(s)


# ── SLIDE 26 ── Content Authenticity Deep Dive ────────────────────────────────
def slide_26():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_BLUE)
    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "26  |  COMPETITOR ANALYSIS", font_size=14, color=ACCENT_BLUE, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Content Authenticity  (\u201cIs this real?\u201d)", font_size=32, color=WHITE, bold=True)
    text_box(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "These systems validate whether content is AI-generated or traceable",
             font_size=14, color=LIGHT_GREY)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_BLUE)

    comps = [
        ("Adobe Content Credentials", "Cryptographic metadata for media provenance", ACCENT_BLUE, [
            ("Focus", "Cryptographic metadata for\nmedia provenance"),
            ("Strength", "Strong push toward industry standard\nfor content attribution"),
            ("Limitation", "Does not tie content back to a\npersistent identity or AI persona"),
        ]),
        ("Google SynthID", "Invisible watermarking of AI content", ACCENT_PURPLE, [
            ("Focus", "Invisible watermarking of AI-generated\ntext / images / audio"),
            ("Strength", "Scalable AI detection and\nattribution layer"),
            ("Limitation", "Detection-focused, not ownership or\neconomic rights enforcement"),
        ]),
    ]
    for ci, (name, subtitle, color, details) in enumerate(comps):
        x = Inches(0.8 + ci * 6.2)
        y = Inches(2.3)
        rect(s, x, y, Inches(5.7), Inches(4.0), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Inches(5.7), Pt(4), fill_color=color)
        text_box(s, x + Inches(0.3), y + Inches(0.2), Inches(5.1), Inches(0.35),
                 name, font_size=18, color=WHITE, bold=True)
        text_box(s, x + Inches(0.3), y + Inches(0.6), Inches(5.1), Inches(0.3),
                 subtitle, font_size=11, color=DIM_GREY)
        for j, (label, text) in enumerate(details):
            dy = y + Inches(1.0 + j * 1.0)
            text_box(s, x + Inches(0.3), dy, Inches(1.2), Inches(0.25),
                     label, font_size=10, color=color, bold=True)
            text_box(s, x + Inches(0.3), dy + Inches(0.25), Inches(5.1), Inches(0.65),
                     text, font_size=12, color=LIGHT_GREY)

    rect(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6), fill_color=RGBColor(0x0E, 0x16, 0x28), radius=0.04)
    text_box(s, Inches(1.0), Inches(6.65), Inches(11), Inches(0.5),
             "Key gap: Provenance exists, but ownership + monetization rights are missing.",
             font_size=14, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_BLUE)
    page_number(s)


# ── SLIDE 27 ── AI Avatar Platforms Deep Dive ─────────────────────────────────
def slide_27():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_PURPLE)
    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "27  |  COMPETITOR ANALYSIS", font_size=14, color=ACCENT_PURPLE, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "AI Avatar / Synthetic Identity Platforms  (\u201cWhat do you look like?\u201d)", font_size=28, color=WHITE, bold=True)
    text_box(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "These tools create digital representations but are disconnected from identity & ownership",
             font_size=13, color=LIGHT_GREY)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_PURPLE)

    comps = [
        (Inches(0.8), "HeyGen", "AI-generated talking avatars for video content", ACCENT_TEAL, [
            ("Strength", "High-quality avatar synthesis\nfor enterprise content creation"),
            ("Limitation", "Avatars are not persistent\nidentity-linked assets"),
        ]),
        (Inches(5.0), "Synthesia", "Enterprise video generation using AI presenters", ACCENT_BLUE, [
            ("Strength", "Scalable corporate training\nand communication content"),
            ("Limitation", "No persistent \u201cdigital self\u201d layer\nacross platforms"),
        ]),
        (Inches(9.2), "D-ID", "Photo-to-video talking avatars", ACCENT_GREEN, [
            ("Strength", "Lightweight avatar generation\nfrom static images"),
            ("Limitation", "Lacks identity binding, commerce,\nor behavioral memory"),
        ]),
    ]
    for (x, name, subtitle, color, details) in comps:
        y = Inches(2.3)
        rect(s, x, y, Inches(3.5), Inches(3.8), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Inches(3.5), Pt(4), fill_color=color)
        text_box(s, x + Inches(0.2), y + Inches(0.2), Inches(3.1), Inches(0.35),
                 name, font_size=16, color=WHITE, bold=True)
        text_box(s, x + Inches(0.2), y + Inches(0.6), Inches(3.1), Inches(0.5),
                 subtitle, font_size=10, color=DIM_GREY)
        for j, (label, text) in enumerate(details):
            dy = y + Inches(1.3 + j * 1.2)
            text_box(s, x + Inches(0.2), dy, Inches(1.0), Inches(0.2),
                     label, font_size=9, color=color, bold=True)
            text_box(s, x + Inches(0.2), dy + Inches(0.22), Inches(3.1), Inches(0.7),
                     text, font_size=11, color=LIGHT_GREY)

    rect(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.6), fill_color=RGBColor(0x1C, 0x14, 0x28), radius=0.04)
    text_box(s, Inches(1.0), Inches(6.55), Inches(11), Inches(0.5),
             "Key gap: Avatars are synthetic outputs, not owned digital entities.",
             font_size=14, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_PURPLE)
    page_number(s)


# ── SLIDE 28 ── AI Companions Deep Dive ───────────────────────────────────────
def slide_28():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_GREEN)
    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "28  |  COMPETITOR ANALYSIS", font_size=14, color=ACCENT_GREEN, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "AI Companions  (\u201cWho interacts with you?\u201d)", font_size=32, color=WHITE, bold=True)
    text_box(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "These platforms focus on conversational personality, not identity or real-world integration",
             font_size=14, color=LIGHT_GREY)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_GREEN)

    comps = [
        ("Replika", "Emotional AI companionship", ACCENT_GREEN, [
            ("Focus", "Emotional AI companionship and\nrelationship simulation"),
            ("Strength", "Deep personalization and\nmemory-driven interactions"),
            ("Limitation", "No verified identity, no real-world\nownership or commerce integration"),
        ]),
        ("Character AI", "User-generated AI characters & roleplay", ACCENT_BLUE, [
            ("Focus", "User-generated AI characters\nand roleplay experiences"),
            ("Strength", "Massive creative ecosystem\nof personalities"),
            ("Limitation", "No persistent identity layer or\neconomic ownership model"),
        ]),
    ]
    for ci, (name, subtitle, color, details) in enumerate(comps):
        x = Inches(0.8 + ci * 6.2)
        y = Inches(2.3)
        rect(s, x, y, Inches(5.7), Inches(4.0), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Inches(5.7), Pt(4), fill_color=color)
        text_box(s, x + Inches(0.3), y + Inches(0.2), Inches(5.1), Inches(0.35),
                 name, font_size=18, color=WHITE, bold=True)
        text_box(s, x + Inches(0.3), y + Inches(0.6), Inches(5.1), Inches(0.3),
                 subtitle, font_size=11, color=DIM_GREY)
        for j, (label, text) in enumerate(details):
            dy = y + Inches(1.0 + j * 1.0)
            text_box(s, x + Inches(0.3), dy, Inches(1.2), Inches(0.25),
                     label, font_size=10, color=color, bold=True)
            text_box(s, x + Inches(0.3), dy + Inches(0.25), Inches(5.1), Inches(0.65),
                     text, font_size=12, color=LIGHT_GREY)

    rect(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6), fill_color=RGBColor(0x0E, 0x1C, 0x14), radius=0.04)
    text_box(s, Inches(1.0), Inches(6.65), Inches(11), Inches(0.5),
             "Key gap: Interaction exists, but it is not anchored to identity, authenticity, or value creation.",
             font_size=14, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_GREEN)
    page_number(s)


# ── SLIDE 29 ── Core Insight — The Missing Layer ──────────────────────────────
def slide_29():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_ORANGE)
    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "29  |  KEY INSIGHT", font_size=14, color=ACCENT_ORANGE, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Core Insight  \u2014  The Missing Layer", font_size=32, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_ORANGE)

    # 4 dimension boxes
    dims = [
        ("\U0001F511", "Identity Systems", "verify who\nyou are", ACCENT_TEAL),
        ("\u2705", "Authenticity Systems", "verify what\nis real", ACCENT_BLUE),
        ("\U0001F5E3\ufe0f", "Avatar Systems", "generate what\nyou look like", ACCENT_PURPLE),
        ("\U0001F9D1\u200D\U0001F4BB", "AI Companions", "simulate how\nyou talk", ACCENT_GREEN),
    ]
    for i, (icon, title, desc, color) in enumerate(dims):
        x = Inches(0.8 + i * 3.1)
        y = Inches(2.3)
        rect(s, x, y, Inches(2.8), Inches(2.5), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Inches(2.8), Pt(4), fill_color=color)
        text_box(s, x + Inches(0.2), y + Inches(0.2), Inches(2.4), Inches(0.5),
                 icon, font_size=26, color=color, alignment=PP_ALIGN.CENTER,
                 font_name="Segoe UI Symbol")
        text_box(s, x + Inches(0.2), y + Inches(0.8), Inches(2.4), Inches(0.35),
                 title, font_size=12, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        text_box(s, x + Inches(0.2), y + Inches(1.3), Inches(2.4), Inches(1.0),
                 desc, font_size=14, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

    # But none combine
    rect(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.6), fill_color=RGBColor(0x1C, 0x14, 0x0E), radius=0.04)
    rect(s, Inches(0.8), Inches(5.2), Pt(5), Inches(1.6), fill_color=ACCENT_ORANGE)
    text_box(s, Inches(1.5), Inches(5.3), Inches(10.5), Inches(0.5),
             "But none combine the full stack:", font_size=15, color=WHITE, bold=True)
    text_box(s, Inches(1.5), Inches(5.8), Inches(10.5), Inches(0.8),
             "Identity  +  Ownership  +  AI Twin  +  Content Provenance  +  Commerce",
             font_size=20, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_ORANGE)
    page_number(s)


# ── SLIDE 30 ── Strategic Implication for MYTAR ──────────────────────────────
def slide_30():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_ORANGE)
    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "30  |  STRATEGIC IMPLICATION", font_size=14, color=ACCENT_ORANGE, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Strategic Implication for MYTAR", font_size=32, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_ORANGE)

    # Left: the gap
    rect(s, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5), fill_color=BG_MID, radius=0.05)
    rect(s, Inches(0.8), Inches(2.3), Inches(5.5), Pt(4), fill_color=ACCENT_RED)
    text_box(s, Inches(1.1), Inches(2.5), Inches(5), Inches(0.35),
             "The Fragmentation Gap", font_size=16, color=ACCENT_RED, bold=True)
    gaps = [
        "No persistent \u201cdigital self\u201d layer",
        "No unified AI-native identity graph",
        "No ownership framework for AI-generated actions/content",
        "No monetization layer for AI twins tied to real identity",
    ]
    for j, g in enumerate(gaps):
        text_box(s, Inches(1.1), Inches(3.1 + j * 0.55), Inches(5), Inches(0.5),
                 f"\u2716  {g}", font_size=13, color=LIGHT_GREY)

    # Right: the white space
    rect(s, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.5), fill_color=RGBColor(0x0E, 0x1C, 0x18), radius=0.05)
    rect(s, Inches(7.0), Inches(2.3), Inches(5.5), Pt(4), fill_color=ACCENT_TEAL)
    text_box(s, Inches(7.3), Inches(2.5), Inches(5), Inches(0.35),
             "The Whitespace MYTAR Targets", font_size=16, color=ACCENT_TEAL, bold=True)
    whitespace = [
        "\u2713  Identity is verifiable",
        "\u2713  AI twins are persistent",
        "\u2713  Content is attributable",
        "\u2713  Actions are ownable",
        "\u2713  Commerce is native to the\ndigital self",
    ]
    for j, w in enumerate(whitespace):
        text_box(s, Inches(7.3), Inches(3.1 + j * 0.55), Inches(5), Inches(0.5),
                 w, font_size=13, color=ACCENT_TEAL, bold=True)

    # Bottom
    text_box(s, Inches(0.8), Inches(6.9), Inches(11), Inches(0.4),
             "This is the structural whitespace MYTAR targets: A unified system for the AI era.",
             font_size=13, color=DIM_GREY, alignment=PP_ALIGN.CENTER)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_ORANGE)
    page_number(s)


# ── SLIDE 31 ── Technology Stack ─────────────────────────────────────────────
def slide_31():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_TEAL)

    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "31  |  TECHNOLOGY", font_size=14, color=ACCENT_TEAL, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Technology Stack", font_size=32, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_TEAL)

    stacks = [
        ("\U0001F3D7\ufe0f", "Blockchain Layer",
         "DID standards (W3C compliant)\nEthereum / Layer-2 chains\nNFT-based identity contracts\nSmart contract royalties", ACCENT_GREEN),
        ("\U0001F9E0", "AI Layer",
         "Large Language Models (OpenAI / open-source)\nVector databases for memory\nVoice synthesis systems\nVideo generation pipelines", ACCENT_PURPLE),
        ("\u2699\ufe0f", "Backend Infrastructure",
         "Microservices architecture\nEvent-driven processing\nSecure API gateway\nScalable cloud infrastructure", ACCENT_BLUE),
        ("\U0001F4BE", "Storage Layer",
         "Encrypted cloud storage\nOptional decentralized storage\nIPFS integration\nData redundancy & backup", ACCENT_ORANGE),
    ]
    for i, (icon, title, content, color) in enumerate(stacks):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(2.3 + (i // 2) * 2.5)
        rect(s, x, y, Inches(5.7), Inches(2.2), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Pt(5), Inches(2.2), fill_color=color)
        text_box(s, x + Inches(0.3), y + Inches(0.15), Inches(0.5), Inches(0.5),
                 icon, font_size=22, color=color, font_name="Segoe UI Symbol")
        text_box(s, x + Inches(0.9), y + Inches(0.15), Inches(4.5), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True)
        text_box(s, x + Inches(0.9), y + Inches(0.6), Inches(4.5), Inches(1.4),
                 content, font_size=11, color=LIGHT_GREY)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_TEAL)
    page_number(s)


# ── SLIDE 32 ── Security & Trust Model ───────────────────────────────────────
def slide_32():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_RED)

    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "32  |  SECURITY", font_size=14, color=ACCENT_RED, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Security & Trust Model", font_size=32, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_RED)

    trust_pillars = [
        ("\U0001F510", "Cryptographic Identity\nSignatures", "Every action signed with\nyour private key"),
        ("\U0001F4CB", "Blockchain Audit\nLogs", "Immutable record of\nall identity events"),
        ("\U0001F6E1\ufe0f", "User-Controlled\nPermissions", "Granular access control\nfor all data"),
        ("\u2705", "Consent-Based\nAvatar Usage", "Explicit permission\nfor every use case"),
        ("\U0001F4DC", "Immutable Interaction\nHistory", "Complete traceability\nof all interactions"),
    ]
    for i, (icon, title, desc) in enumerate(trust_pillars):
        x = Inches(0.8 + (i % 3) * 4.1)
        y = Inches(2.3 + (i // 3) * 2.3)
        rect(s, x, y, Inches(3.7), Inches(2.0), fill_color=BG_MID, radius=0.04)
        rect(s, x, y, Pt(4), Inches(2.0), fill_color=ACCENT_RED)
        text_box(s, x + Inches(0.3), y + Inches(0.2), Inches(0.5), Inches(0.5),
                 icon, font_size=22, color=ACCENT_RED, font_name="Segoe UI Symbol")
        text_box(s, x + Inches(0.9), y + Inches(0.15), Inches(2.5), Inches(0.5),
                 title, font_size=14, color=WHITE, bold=True)
        text_box(s, x + Inches(0.9), y + Inches(0.75), Inches(2.5), Inches(1.0),
                 desc, font_size=11, color=LIGHT_GREY)

    # Bottom
    rect(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.6), fill_color=RGBColor(0x1A, 0x0E, 0x0E), radius=0.04)
    text_box(s, Inches(1.0), Inches(6.55), Inches(11), Inches(0.5),
             "Every digital action can be verified or revoked.",
             font_size=15, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_RED)
    page_number(s)


# ── SLIDE 33 ── Go-To-Market Strategy ────────────────────────────────────────
def slide_33():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_TEAL)

    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "33  |  STRATEGY", font_size=14, color=ACCENT_TEAL, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Go-To-Market Strategy", font_size=32, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_TEAL)

    phases = [
        ("PHASE 1", "High-Trust Individuals", ACCENT_GREEN,
         "CEOs  \u2022  Celebrities  \u2022  Politicians  \u2022  Influencers",
         "Identity protection\nDeepfake prevention\nVerified digital presence"),
        ("PHASE 2", "Creator Economy", ACCENT_BLUE,
         "Content creators  \u2022  Artists  \u2022  Influencers",
         "Avatar monetization tools\nLicensing marketplace\nContent generation ecosystem"),
        ("PHASE 3", "Enterprise & Government", ACCENT_PURPLE,
         "Corporations  \u2022  Government agencies  \u2022  Institutions",
         "Identity verification APIs\nFraud prevention systems\nSecure communication infra"),
        ("PHASE 4", "Consumer Scale", ACCENT_ORANGE,
         "Global consumers  \u2022  Retail  \u2022  E-commerce",
         "Fashion AI\nShopping assistant\nPersonal digital twins"),
    ]
    for i, (phase, title, color, audience, focus) in enumerate(phases):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(2.3 + (i // 2) * 2.4)
        rect(s, x, y, Inches(5.7), Inches(2.1), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Pt(5), Inches(2.1), fill_color=color)
        text_box(s, x + Inches(0.3), y + Inches(0.1), Inches(2), Inches(0.3),
                 phase, font_size=11, color=color, bold=True)
        text_box(s, x + Inches(2.3), y + Inches(0.1), Inches(3), Inches(0.3),
                 title, font_size=16, color=WHITE, bold=True)
        text_box(s, x + Inches(0.3), y + Inches(0.5), Inches(5.1), Inches(0.3),
                 audience, font_size=10, color=DIM_GREY)
        text_box(s, x + Inches(0.3), y + Inches(0.9), Inches(5.1), Inches(1.0),
                 focus, font_size=12, color=LIGHT_GREY)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_TEAL)
    page_number(s)


# ── SLIDE 34 ── Business Model ───────────────────────────────────────────────
def slide_34():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_GREEN)

    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "34  |  BUSINESS", font_size=14, color=ACCENT_GREEN, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Business Model", font_size=32, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_GREEN)

    revenues = [
        ("\U0001F4B3", "Identity Subscription Plans", "Recurring revenue from\nverified identity hosting", ACCENT_TEAL),
        ("\U0001F3F7", "NFT Identity Issuance Fees", "One-time minting fee for\nidentity certificates", ACCENT_BLUE),
        ("\U0001F3AC", "Avatar Usage Fees", "Per-use or subscription for\nAI twin interactions", ACCENT_PURPLE),
        ("\U0001F3ED", "Licensing Marketplace Commissions", "Percentage of every\navatar licensing deal", ACCENT_ORANGE),
        ("\U0001F4E1", "Enterprise API Contracts", "B2B identity verification\nand twin infrastructure", ACCENT_GREEN),
        ("\U0001F4B8", "Ecosystem Transaction Fees", "Micro-fees on every\nidentity-authenticated action", ACCENT_RED),
    ]
    for i, (icon, title, desc, color) in enumerate(revenues):
        x = Inches(0.8 + (i % 3) * 4.1)
        y = Inches(2.3 + (i // 3) * 2.3)
        rect(s, x, y, Inches(3.7), Inches(2.0), fill_color=BG_MID, radius=0.04)
        rect(s, x, y, Pt(4), Inches(2.0), fill_color=color)
        text_box(s, x + Inches(0.3), y + Inches(0.2), Inches(0.5), Inches(0.5),
                 icon, font_size=22, color=color, font_name="Segoe UI Symbol")
        text_box(s, x + Inches(0.9), y + Inches(0.15), Inches(2.5), Inches(0.5),
                 title, font_size=14, color=WHITE, bold=True)
        text_box(s, x + Inches(0.9), y + Inches(0.75), Inches(2.5), Inches(1.0),
                 desc, font_size=11, color=LIGHT_GREY)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_GREEN)
    page_number(s)


# ── SLIDE 35 ── Strategic Advantage ──────────────────────────────────────────
def slide_35():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_PURPLE)

    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "35  |  ADVANTAGE", font_size=14, color=ACCENT_PURPLE, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Strategic Advantage", font_size=32, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_PURPLE)

    advantages = [
        ("\U0001F3D7\ufe0f", "Blockchain-Native\nIdentity Ownership", "Identity is not rented \u2014\nit is owned by the user"),
        ("\U0001F9E0", "Persistent AI\nDigital Twins", "Continuously evolving\nrepresentation of self"),
        ("\U0001F310", "Unified Identity + Commerce\n+ Communication Layer", "Everything in one\nintegrated platform"),
        ("\U0001F3C6", "Creator Economy\nEmbedded in Identity", "Monetization is built\ninto the identity itself"),
    ]
    for i, (icon, title, desc) in enumerate(advantages):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(2.3 + (i // 2) * 2.3)
        rect(s, x, y, Inches(5.7), Inches(2.0), fill_color=BG_MID, radius=0.05)
        rect(s, x, y, Pt(5), Inches(2.0), fill_color=ACCENT_PURPLE)
        text_box(s, x + Inches(0.3), y + Inches(0.3), Inches(0.6), Inches(0.6),
                 icon, font_size=28, color=ACCENT_PURPLE, font_name="Segoe UI Symbol")
        text_box(s, x + Inches(1.1), y + Inches(0.2), Inches(4.3), Inches(0.5),
                 title, font_size=16, color=WHITE, bold=True)
        text_box(s, x + Inches(1.1), y + Inches(0.9), Inches(4.3), Inches(0.8),
                 desc, font_size=13, color=LIGHT_GREY)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_PURPLE)
    page_number(s)


# ── SLIDE 36 ── End State Vision ─────────────────────────────────────────────
def slide_36():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_TEAL)

    # Large circle decoration
    circle(s, W - Inches(1.5), -Inches(1), Inches(5), fill_color=RGBColor(0x12, 0x1C, 0x30))
    circle(s, -Inches(1), H - Inches(2), Inches(4), fill_color=RGBColor(0x0E, 0x16, 0x26))

    text_box(s, Inches(0.8), Inches(0.2), Inches(2), Inches(0.5),
             "36", font_size=14, color=ACCENT_TEAL, bold=True)

    text_box(s, Inches(0.8), Inches(1.5), Inches(11), Inches(1.2),
             "End State Vision", font_size=48, color=WHITE, bold=True)

    rect(s, Inches(0.8), Inches(2.8), Inches(1.5), Pt(4), fill_color=ACCENT_TEAL)

    text_box(s, Inches(0.8), Inches(3.2), Inches(11), Inches(1.5),
             "MYTAR evolves into:\nThe global protocol for digital human identity, representation,\nand interaction in the AI-powered internet.",
             font_size=22, color=LIGHT_GREY)

    rect(s, Inches(0.8), Inches(4.8), Inches(0.8), Pt(3), fill_color=ACCENT_TEAL)

    visions = [
        "Every human has a verified digital twin",
        "Every interaction is authenticated",
        "Every digital action is traceable and owned",
    ]
    for i, v in enumerate(visions):
        text_box(s, Inches(0.8), Inches(5.1 + i * 0.5), Inches(11), Inches(0.4),
                 f"\u2B9E  {v}", font_size=16, color=ACCENT_TEAL)

    # Bottom tag
    text_box(s, Inches(0.8), Inches(6.8), Inches(6), Inches(0.4),
             "MYTAR  \u2022  The Ownership Layer for Human Identity",
             font_size=12, color=DIM_GREY)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_TEAL)
    page_number(s)


# ── AGENDA SLIDE ───────────────────────────────────────────────────────────────
def slide_agenda():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_TEAL)
    text_box(s, Inches(0.8), Inches(0.2), Inches(4), Inches(0.5),
             "AGENDA", font_size=14, color=ACCENT_TEAL, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "Presentation Outline", font_size=32, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.3), Inches(1.5), Pt(3), fill_color=ACCENT_TEAL)

    agenda_items = [
        ("A", "Context & Vision", "The problem, thesis, and founding belief behind MYTAR", ACCENT_TEAL),
        ("B", "Platform Architecture", "Identity layer, blockchain core, NFT ownership, content authenticity", ACCENT_BLUE),
        ("C", "Digital Twin System", "AI twin engine, brain layer, and multi-mode interaction", ACCENT_PURPLE),
        ("D", "Use Case Ecosystem", "Fashion, commerce, proxy, communication, creator economy, legacy", ACCENT_GREEN),
        ("E", "Competitive Landscape", "Market fragmentation, competitor deep-dives, strategic whitespace", ACCENT_ORANGE),
        ("F", "Technology & Security", "Tech stack, security model, cryptographic trust infrastructure", ACCENT_RED),
        ("G", "Go-To-Market & Business", "GTM phases, revenue streams, strategic advantage, end-state vision", ACCENT_TEAL),
    ]
    for i, (letter, title, desc, color) in enumerate(agenda_items):
        y = Inches(1.7 + i * 0.7)
        circle(s, Inches(0.8), y + Inches(0.05), Inches(0.4), fill_color=color)
        text_box(s, Inches(0.8), y + Inches(0.05), Inches(0.4), Inches(0.4),
                 letter, font_size=14, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        text_box(s, Inches(1.5), y, Inches(3.5), Inches(0.3),
                 title, font_size=15, color=WHITE, bold=True)
        text_box(s, Inches(1.5), y + Inches(0.3), Inches(10), Inches(0.3),
                 desc, font_size=10, color=LIGHT_GREY)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_TEAL)
    page_number(s)


# ── SECTION DIVIDER HELPER ─────────────────────────────────────────────────────
def section_divider(section_letter, section_title, section_subtitle, accent_color):
    s = add_blank_slide()
    bg(s, BG_DARK)
    # Full-width accent block
    rect(s, 0, 0, W, H, fill_color=RGBColor(0x0E, 0x14, 0x24))
    rect(s, 0, 0, W, Pt(5), fill_color=accent_color)
    # Large section letter
    text_box(s, Inches(1.0), Inches(1.5), Inches(3), Inches(1.5),
             section_letter, font_size=96, color=accent_color, bold=True, font_name="Calibri")
    # Section title
    text_box(s, Inches(1.0), Inches(3.2), Inches(10), Inches(0.8),
             section_title, font_size=36, color=WHITE, bold=True)
    # Subtitle
    text_box(s, Inches(1.0), Inches(4.0), Inches(10), Inches(0.5),
             section_subtitle, font_size=16, color=LIGHT_GREY)
    # Accent line
    rect(s, Inches(1.0), Inches(4.7), Inches(2.0), Pt(3), fill_color=accent_color)
    # Right decorative element
    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=accent_color)
    page_number(s)


# ── BUILD PRESENTATION ────────────────────────────────────────────────────────
# Actually slide_05 is wrong - it makes a blank slide we don't use. Let's fix that.
# We need to restructure: slide_05 creates s first, then s2. The `s` is unused.

# Clear accidentally created slides
# We'll just call functions in order. The slide_05 currently creates 2 slides.
# Let's fix by making slide_05 not create the first one.

# Re-defining slide_05 properly (the version above creates a second - let's just
# call these properly)

# Actually let me just rewrite slide_05 as a clean function and not use the first s.

def slide_05_fixed():
    s = add_blank_slide()
    bg(s, BG_DARK)
    rect(s, 0, 0, W, Pt(3), fill_color=ACCENT_TEAL)
    text_box(s, Inches(0.8), Inches(0.2), Inches(2), Inches(0.5),
             "05", font_size=14, color=ACCENT_TEAL, bold=True)
    text_box(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
             "MYTAR System Overview", font_size=32, color=WHITE, bold=True)
    text_box(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "A 4-layer ecosystem built on cryptographic trust",
             font_size=14, color=LIGHT_GREY)
    rect(s, Inches(0.8), Inches(1.8), Inches(1.5), Pt(3), fill_color=ACCENT_TEAL)

    # ── ARCHITECTURE DIAGRAM ──
    layers = [
        ("APPLICATION LAYER", "How you interact with the world", ACCENT_TEAL,
         ["AI Fashion & Commerce", "Personal AI Shopper", "Communication Layer",
          "AR/VR Interfaces", "Meeting Proxy", "Creator Marketplace"]),
        ("AI TWIN LAYER", "How you exist digitally", ACCENT_BLUE,
         ["Voice & Face Models", "Memory Graph", "Behavioral Engine",
          "Knowledge Base", "Context Reasoning", "Multi-modal Output"]),
        ("OWNERSHIP LAYER", "What you own", ACCENT_PURPLE,
         ["NFT Identity Certificate", "Digital Rights Contract",
          "Licensing Engine", "Royalty Distribution", "Usage Audit"]),
        ("IDENTITY LAYER \u2014 Blockchain Core", "Who you are", ACCENT_GREEN,
         ["DID Standards", "Identity Registry", "Immutable Proofs",
          "Crypto Wallet", "Content Signatures"]),
    ]

    left_box = Inches(0.8)
    box_w = Inches(5.5)
    box_h = Inches(1.15)
    gap = Inches(0.15)
    start_y = Inches(2.3)
    detail_x = Inches(7.0)
    detail_w = Inches(5.8)

    for i, (title, subtitle, color, items) in enumerate(layers):
        y = start_y + i * (box_h + gap)
        shape = rect(s, left_box, y, box_w, box_h, fill_color=RGBColor(0x12, 0x1C, 0x38), radius=0.04)
        rect(s, left_box, y, Pt(5), box_h, fill_color=color)
        circle(s, left_box + Inches(0.25), y + Inches(0.15), Inches(0.35), fill_color=color)
        text_box(s, left_box + Inches(0.25), y + Inches(0.15), Inches(0.35), Inches(0.35),
                 str(4 - i), font_size=14, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        text_box(s, left_box + Inches(0.75), y + Inches(0.1), Inches(4.5), Inches(0.35),
                 title, font_size=14, color=color, bold=True)
        text_box(s, left_box + Inches(0.75), y + Inches(0.5), Inches(4.5), Inches(0.3),
                 subtitle, font_size=11, color=LIGHT_GREY)

        if i < 3:
            arrow_y = y + box_h
            arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left_box + Inches(0.5), arrow_y,
                                        Inches(0.3), gap)
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = color
            arrow.line.fill.background()

        item_y = y + Inches(0.05)
        for j, item in enumerate(items):
            ix = detail_x + (j % 2) * (detail_w / 2)
            iy = item_y + (j // 2) * Inches(0.32)
            text_box(s, ix, iy, detail_w / 2 - Inches(0.1), Inches(0.3),
                     f"\u25B8 {item}", font_size=10, color=LIGHT_GREY)

    text_box(s, left_box, Inches(6.9), Inches(8), Inches(0.4),
             "Each layer builds on cryptographic trust . Identity  \u2192  Ownership  \u2192  AI Twin  \u2192  Applications",
             font_size=12, color=DIM_GREY)

    rect(s, W - Inches(0.1), 0, Pt(3), H, fill_color=ACCENT_TEAL)
    page_number(s)
    return s


# ── Generate all slides in order ─────────────────────────────────────────────

def generate():
    # Opening
    slide_01()                       # 1  — Title
    slide_agenda()                   # 2  — Agenda

    # Section A: Context & Vision
    section_divider("A", "Context & Vision", "The problem, thesis, and founding belief behind MYTAR", ACCENT_TEAL)  # 3
    slide_02()                       # 4  — Founding Vision
    slide_03()                       # 5  — Core Thesis
    slide_04()                       # 6  — Macro Problem

    # Section B: Platform Architecture
    section_divider("B", "Platform Architecture", "Identity layer, blockchain core, NFT ownership, content authenticity", ACCENT_BLUE)  # 7
    slide_05_fixed()                 # 8  — System Overview
    slide_06()                       # 9  — Identity Layer
    slide_07()                       # 10 — Identity Verification
    slide_08()                       # 11 — NFT Ownership
    slide_09()                       # 12 — Blockchain Utility
    slide_10()                       # 13 — Content Authenticity

    # Section C: Digital Twin System
    section_divider("C", "Digital Twin System", "AI twin engine, brain layer, and multi-mode interaction", ACCENT_PURPLE)  # 14
    slide_11()                       # 15 — Digital Twin Engine
    slide_12()                       # 16 — AI Brain Layer
    slide_13()                       # 17 — Interaction Modes

    # Section D: Use Case Ecosystem
    section_divider("D", "Use Case Ecosystem", "Fashion, commerce, proxy, communication, creator economy, and legacy", ACCENT_GREEN)  # 18
    slide_14()                       # 19 — AI Fashion
    slide_15()                       # 20 — Personal AI Shopper
    slide_16()                       # 21 — Body & Appearance
    slide_17()                       # 22 — Digital Proxy
    slide_18()                       # 23 — Global Communication
    slide_19()                       # 24 — Creator Economy
    slide_20()                       # 25 — Digital Legacy
    slide_21()                       # 26 — Share Your Tar
    slide_22()                       # 27 — Future Interfaces

    # Section E: Competitive Landscape
    section_divider("E", "Competitive Landscape", "Market fragmentation, competitor deep-dives, and strategic whitespace", ACCENT_ORANGE)  # 28
    slide_23()                       # 29 — Competitive Overview
    slide_24()                       # 30 — Market Fragmentation
    slide_25()                       # 31 — Identity Systems
    slide_26()                       # 32 — Content Authenticity
    slide_27()                       # 33 — AI Avatars
    slide_28()                       # 34 — AI Companions
    slide_29()                       # 35 — Core Insight
    slide_30()                       # 36 — Strategic Implication

    # Section F: Technology & Security
    section_divider("F", "Technology & Security", "Tech stack, security model, and cryptographic trust infrastructure", ACCENT_RED)  # 37
    slide_31()                       # 38 — Technology Stack
    slide_32()                       # 39 — Security & Trust

    # Section G: Go-To-Market & Business
    section_divider("G", "Go-To-Market & Business", "GTM phases, revenue streams, strategic advantage, and end-state vision", ACCENT_TEAL)  # 40
    slide_33()                       # 41 — GTM Strategy
    slide_34()                       # 42 — Business Model
    slide_35()                       # 43 — Strategic Advantage
    slide_36()                       # 44 — End State Vision

    output_path = "/Users/mohitmendiratta/Projects/agents/job1/MYTAR_Investor_Deck_v3.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    generate()
