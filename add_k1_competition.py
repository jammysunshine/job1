from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import re

INPUT  = "/Users/mohitmendiratta/Projects/agents/job1/MYTAR-K1.pptx"
OUTPUT = "/Users/mohitmendiratta/Projects/agents/job1/MYTAR-K1.pptx"

W = 16256000
H = 9144000
FONT = "Liter"

# Colors
GOLD    = RGBColor(0xC8, 0xA5, 0x5A)
WHITE   = RGBColor(0xF0, 0xF0, 0xF0)
GREY    = RGBColor(0x8A, 0x8D, 0x93)
CARD_BG = RGBColor(0x15, 0x16, 0x1C)
BG      = RGBColor(0x0B, 0x0C, 0x10)


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[0])


def rect(slide, left, top, width, height, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def rect_rnd(slide, left, top, width, height, fill, rad=0.02):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = rad
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def txt(slide, left, top, width, height, text, fs=Pt(11), color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = fs
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    return tb


def header(slide, title, subtitle=None):
    txt(slide, 762000, 355600, 1016000, 228600, "MYTAR", fs=Pt(11), color=GREY)
    txt(slide, 762000, 762000, 12000000, 571500, title, fs=Pt(30), color=GOLD)
    rect(slide, 762000, 1422400, 762000, 38100, GOLD)
    if subtitle:
        txt(slide, 762000, 1587500, 12000000, 355600, subtitle, fs=Pt(16), color=GREY)


def bottom_bar(slide, section_label):
    rect(slide, 762000, 8763000, 14732000, 12700, GOLD)
    txt(slide, 11500000, 8445500, 3500000, 228600, section_label, fs=Pt(11), color=GREY, align=PP_ALIGN.RIGHT)


def section_05_label():
    return "05 / COMPETITIVE LANDSCAPE"


#
# Slide 31: Market Fragmentation
#
def slide_market_fragmentation(prs):
    s = add_slide(prs)
    header(s, "MARKET FRAGMENTATION",
           "The market is split across four emerging but disconnected categories")

    cats = [
        ("IDENTITY SYSTEMS", "\u2192 Who are you?", [
            "Identity verification,\ncredentials, anti-bot proofs"]),
        ("CONTENT AUTHENTICITY", "\u2192 Is this real?", [
            "Provenance tracking,\nAI watermarking, attribution"]),
        ("AI AVATAR PLATFORMS", "\u2192 What do you look like?", [
            "Avatar generation,\nvideo synthesis, enterprise\ntools"]),
        ("AI COMPANIONS", "\u2192 Who interacts with you?", [
            "Emotional AI, roleplay\ncharacters, personal\nassistants"]),
    ]

    cx = 762000
    cy = 2100000
    cw = 3429000
    ch = 3700000
    gap = 254000

    for i, (title, q, bullets) in enumerate(cats):
        x = cx + i * (cw + gap)
        rect_rnd(s, x, cy, cw, ch, CARD_BG)
        txt(s, x + 190500, cy + 254000, cw - 381000, 274320, title, fs=Pt(11), color=GOLD)
        txt(s, x + 190500, cy + 650000, cw - 381000, 274320, q, fs=Pt(15), color=WHITE)
        for j, b in enumerate(bullets):
            txt(s, x + 190500, cy + 1050000 + j * 450000, cw - 381000, 400000,
                b, fs=Pt(13), color=GREY)

    # Bottom insight bar
    bx, by = 762000, 6100000
    rect(s, bx, by, 14732000, 1900000, CARD_BG)
    rect(s, bx, by, 14732000, 50800, GOLD)
    txt(s, 1200000, by + 380000, 6000000, 279400, "THE FRAGMENTATION PROBLEM", fs=Pt(14), color=GOLD)
    txt(s, 1200000, by + 800000, 13500000, 800000,
        "Despite rapid innovation in digital identity, AI avatars, and content authenticity, "
        "today\u2019s ecosystem is split into isolated vertical stacks, "
        "each solving only a narrow part of the problem.",
        fs=Pt(18), color=WHITE)

    bottom_bar(s, section_05_label())
    return len(prs.slides) - 1


#
# Slide: Competitor Deep Dive (2 competitors)
#
def slide_competitor_2col(prs, title, subtitle, competitors, accent=None):
    """
    competitors = [(name, subtitle?, [(label, value), ...]), ...]
    """
    s = add_slide(prs)
    header(s, title, subtitle)

    cy = 2100000
    ch = 4200000
    cw = 7239000
    gap = 254000

    for ci, (name, sub, details) in enumerate(competitors):
        x = 762000 + ci * (cw + gap)
        rect_rnd(s, x, cy, cw, ch, CARD_BG)

        # Name
        txt(s, x + 190500, cy + 200000, cw - 381000, 350000, name.upper(), fs=Pt(20), color=GOLD)
        if sub:
            txt(s, x + 190500, cy + 580000, cw - 381000, 250000, sub, fs=Pt(11), color=GREY)

        # Detail rows
        row_y = cy + 900000 if sub else cy + 650000
        for j, (label, value) in enumerate(details):
            ry = row_y + j * 950000
            txt(s, x + 190500, ry, 1800000, 228600, label.upper(), fs=Pt(10), color=GOLD)
            txt(s, x + 190500, ry + 280000, cw - 500000, 550000, value, fs=Pt(13), color=WHITE)

    bottom_bar(s, section_05_label())
    return len(prs.slides) - 1


#
# Slide: Competitor Deep Dive (3 competitors) - for AI Avatars
#
def slide_competitor_3col(prs, title, subtitle, competitors):
    s = add_slide(prs)
    header(s, title, subtitle)

    cy = 2100000
    ch = 4200000
    cw = 4572000
    gap = 254000
    x0 = 762000
    total_w = 3 * cw + 2 * gap  # 3*4572000 + 2*254000 = 14224000
    left_start = (W - total_w) // 2

    for ci, (name, details) in enumerate(competitors):
        x = left_start + ci * (cw + gap)
        rect_rnd(s, x, cy, cw, ch, CARD_BG)

        txt(s, x + 190500, cy + 200000, cw - 381000, 350000, name.upper(), fs=Pt(18), color=GOLD, align=PP_ALIGN.CENTER)

        row_y = cy + 700000
        for j, (label, value) in enumerate(details):
            ry = row_y + j * 1050000
            txt(s, x + 190500, ry, cw - 381000, 228600, label.upper(), fs=Pt(10), color=GOLD, align=PP_ALIGN.CENTER)
            txt(s, x + 190500, ry + 280000, cw - 381000, 650000, value, fs=Pt(12), color=WHITE, align=PP_ALIGN.CENTER)

    bottom_bar(s, section_05_label())
    return len(prs.slides) - 1


#
# Slide: Core Insight
#
def slide_core_insight(prs):
    s = add_slide(prs)
    header(s, "CORE INSIGHT \u2014 THE MISSING LAYER",
           "Across all categories, systems are optimized for one dimension only")

    dims = [
        ("IDENTITY SYSTEMS", "\u2192 verify who you are", GOLD),
        ("AUTHENTICITY SYSTEMS", "\u2192 verify what is real", GOLD),
        ("AVATAR SYSTEMS", "\u2192 generate what you look like", GOLD),
        ("AI COMPANIONS", "\u2192 simulate how you talk", GOLD),
    ]

    cx = 762000
    cy = 2200000
    cw = 3429000
    ch = 1600000
    gap = 254000

    for i, (label, desc, color) in enumerate(dims):
        x = cx + i * (cw + gap)
        rect_rnd(s, x, cy, cw, ch, CARD_BG)
        txt(s, x + 190500, cy + 250000, cw - 381000, 274320, label, fs=Pt(11), color=GOLD)
        txt(s, x + 190500, cy + 700000, cw - 381000, 500000, desc, fs=Pt(15), color=WHITE)

    # Bottom: "But none combine the full stack"
    by = 4200000
    rect(s, 762000, by, 14732000, 1500000, CARD_BG)
    rect(s, 762000, by, 14732000, 50800, GOLD)
    txt(s, 1200000, by + 300000, 6000000, 279400, "BUT NONE COMBINE THE FULL STACK", fs=Pt(14), color=GOLD)
    txt(s, 1200000, by + 700000, 13500000, 600000,
        "Identity + Ownership + AI Twin + Content Provenance + Commerce",
        fs=Pt(20), color=WHITE, bold=True)

    bottom_bar(s, section_05_label())
    return len(prs.slides) - 1


#
# Slide: Strategic Implication
#
def slide_strategic_implication(prs):
    s = add_slide(prs)
    header(s, "STRATEGIC IMPLICATION FOR MYTAR",
           "This fragmentation creates a structural gap in the market")

    cy = 2100000
    ch = 4400000
    cw = 7239000
    gap = 254000

    # Left: Fragmentation Gap
    x = 762000
    rect_rnd(s, x, cy, cw, ch, CARD_BG)
    txt(s, x + 190500, cy + 250000, cw - 381000, 350000, "THE FRAGMENTATION GAP", fs=Pt(16), color=GOLD)
    gaps = [
        "No persistent \u201cdigital self\u201d layer",
        "No unified AI-native identity graph",
        "No ownership framework for AI\nactions / content",
        "No monetization layer for AI twins\ntied to real identity",
    ]
    for j, g in enumerate(gaps):
        txt(s, x + 190500, cy + 850000 + j * 700000, cw - 381000, 600000,
            g, fs=Pt(15), color=WHITE)

    # Right: Whitespace
    x = 762000 + cw + gap
    rect_rnd(s, x, cy, cw, ch, CARD_BG)
    txt(s, x + 190500, cy + 250000, cw - 381000, 350000, "THE WHITESPACE MYTAR TARGETS", fs=Pt(16), color=GOLD)
    ws = [
        "Identity is verifiable",
        "AI twins are persistent",
        "Content is attributable",
        "Actions are ownable",
        "Commerce is native to the\ndigital self",
    ]
    for j, w in enumerate(ws):
        txt(s, x + 190500, cy + 850000 + j * 700000, cw - 381000, 600000,
            w, fs=Pt(15), color=WHITE)

    bottom_bar(s, section_05_label())
    return len(prs.slides) - 1


#
# ── Build competitor content ──
#

def build_identity_slide(prs):
    return slide_competitor_2col(prs,
        "IDENTITY SYSTEMS",
        "These platforms verify human identity but do not extend to AI representation or digital ownership",
        [
            ("World ID", "Tools for Humanity", [
                ("Focus", "Proof-of-human via biometric scanning /\niris-based identity"),
                ("Strength", "Strong anti-bot / Sybil resistance\nfor digital platforms"),
                ("Limitation", "No concept of AI persona, content\nownership, or monetization layer"),
            ]),
            ("Microsoft Entra Verified ID", "Enterprise verifiable credentials", [
                ("Focus", "Enterprise-grade verifiable credentials\n(work, education, access)"),
                ("Strength", "Deep enterprise integration,\ncompliance-ready identity verification"),
                ("Limitation", "Static identity layer; not designed for\nconsumer AI or creative identity extension"),
            ]),
        ])


def build_content_slide(prs):
    return slide_competitor_2col(prs,
        "CONTENT AUTHENTICITY",
        "These systems validate whether content is AI-generated or traceable, but remain disconnected from identity and ownership",
        [
            ("Adobe Content Credentials", "Cryptographic metadata for media provenance", [
                ("Focus", "Cryptographic metadata for\nmedia provenance"),
                ("Strength", "Strong push toward industry standard\nfor content attribution"),
                ("Limitation", "Does not tie content back to a\npersistent identity or AI persona"),
            ]),
            ("Google SynthID", "Invisible AI content watermarking", [
                ("Focus", "Invisible watermarking of AI-generated\ntext / images / audio"),
                ("Strength", "Scalable AI detection and\nattribution layer"),
                ("Limitation", "Detection-focused, not ownership or\neconomic rights enforcement"),
            ]),
        ])


def build_avatar_slide(prs):
    return slide_competitor_3col(prs,
        "AI AVATAR / SYNTHETIC IDENTITY PLATFORMS",
        "These tools create digital representations but are disconnected from real identity systems and ownership rails",
        [
            ("HeyGen", [
                ("Focus", "AI-generated talking\navatars for video"),
                ("Strength", "High-quality avatar\nsynthesis for enterprise"),
                ("Limitation", "Avatars not persistent\nidentity-linked assets"),
            ]),
            ("Synthesia", [
                ("Focus", "Enterprise video\ngeneration using AI"),
                ("Strength", "Scalable corporate\ntraining and content"),
                ("Limitation", "No persistent \u201cdigital\nself\u201d layer"),
            ]),
            ("D-ID", [
                ("Focus", "Photo-to-video\ntalking avatars"),
                ("Strength", "Lightweight avatar gen\nfrom static images"),
                ("Limitation", "Lacks identity binding,\ncommerce, or memory"),
            ]),
        ])


def build_companion_slide(prs):
    return slide_competitor_2col(prs,
        "AI COMPANIONS",
        "These platforms focus on conversational personality and engagement, not identity verification or real-world integration",
        [
            ("Replika", "Emotional AI companionship", [
                ("Focus", "Emotional AI companionship and\nrelationship simulation"),
                ("Strength", "Deep personalization and\nmemory-driven interactions"),
                ("Limitation", "No verified identity, no real-world\nownership or commerce integration"),
            ]),
            ("Character AI", "User-generated AI characters", [
                ("Focus", "User-generated AI characters\nand roleplay experiences"),
                ("Strength", "Massive creative ecosystem\nof personalities"),
                ("Limitation", "No persistent identity layer or\neconomic ownership model"),
            ]),
        ])


#
# ── Reorder slides ──
#

def reorder_slides(prs, new_order):
    ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    pres_elem = prs.part._element
    sldIdLst = pres_elem.find(f'{{{ns}}}sldIdLst')
    elems = list(sldIdLst)
    ordered = [elems[i] for i in new_order]
    for e in elems:
        sldIdLst.remove(e)
    for e in ordered:
        sldIdLst.append(e)


#
# ── Main ──
#

def main():
    prs = Presentation(INPUT)
    orig = len(prs.slides)
    print(f"Original: {orig} slides")

    # Only add if not already added (check if slide 31 exists and has competition content)
    already_added = False
    if orig > 30:
        s = prs.slides[30]
        for sh in s.shapes:
            if sh.has_text_frame and "MARKET FRAGMENTATION" in sh.text_frame.text:
                already_added = True
                print("Competition slides already present, skipping.")
                break

    if not already_added:
        print("Creating competition detail slides...")
        new_indices = []
        builders = [
            slide_market_fragmentation,
            build_identity_slide,
            build_content_slide,
            build_avatar_slide,
            build_companion_slide,
            slide_core_insight,
            slide_strategic_implication,
        ]
        for fn in builders:
            idx = fn(prs)
            new_indices.append(idx)

        total = len(prs.slides)
        print(f"Total after append: {total}")
        print(f"New slide indices: {new_indices}")

        # New order: 0-29 (slides 1-30), then new slides, then 30-38 (old slides 31-39)
        # Actually: new slides are at indices 39-45 (orig was 39)
        # New order: [0-29, 39-45, 30-38]
        new_order = list(range(30))  # 0-29 (original 1-30)
        new_order.extend(new_indices)  # new slides
        new_order.extend(range(30, orig))  # original 31-39 (now at positions 38-45)

        print("Reordering...")
        reorder_slides(prs, new_order)

        prs.save(OUTPUT)
        print(f"Saved! Total: {len(prs.slides)}")


if __name__ == "__main__":
    main()
