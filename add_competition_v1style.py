from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

# v1 exact measurements
W = 12191695
H = 6858000
TOP_BAR_H = 38100
RIGHT_BAR_W = 38100
ACCNT_STRIP_H = 50800  # card top strip
LEFT_ACCNT_W = 63500   # left accent bar (wide variant)
LEFT_ACCNT_W2 = 50800  # left accent bar (narrow variant)

# v1 positions
LABEL_X = 731520
LABEL_Y = 182880
LABEL_W = 3657600
LABEL_H = 457200

TITLE_X = 731520
TITLE_Y = 640080
TITLE_W = 9144000
TITLE_H = 548640

SUBTITLE_X = 731520
SUBTITLE_Y = 1188720
SUBTITLE_W = 9144000
SUBTITLE_H = 365760

ACCNT_LINE_X = 731520
ACCNT_LINE_Y = 1645920
ACCNT_LINE_W = 1371600
ACCNT_LINE_H = 38100

PAGE_NUM_X = 10820095
PAGE_NUM_Y = 6400800
PAGE_NUM_W = 1097280
PAGE_NUM_H = 365760

# v1 colors
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_LTGREY   = RGBColor(0xBB, 0xBF, 0xD0)
C_DIMGREY  = RGBColor(0x6B, 0x72, 0x8A)
C_TEAL     = RGBColor(0x33, 0xE0, 0xCC)
C_BLUE     = RGBColor(0x29, 0x80, 0xB9)
C_PURPLE   = RGBColor(0x8B, 0x5C, 0xF6)
C_GREEN    = RGBColor(0x10, 0xB9, 0x81)
C_ORANGE   = RGBColor(0xF5, 0x9E, 0x0B)
C_RED      = RGBColor(0xEF, 0x44, 0x44)
C_CARD     = RGBColor(0x11, 0x16, 0x28)
C_CARD2    = RGBColor(0x12, 0x1C, 0x38)
C_BOTDARK  = RGBColor(0x1C, 0x14, 0x0E)
C_BOTTEAL  = RGBColor(0x0E, 0x1C, 0x18)
C_BOTBLUE  = RGBColor(0x0E, 0x16, 0x28)
C_BOTPURP  = RGBColor(0x1C, 0x14, 0x28)
C_BOTGREEN = RGBColor(0x0E, 0x1C, 0x14)

INPUT  = "/Users/mohitmendiratta/Projects/agents/job1/MYTAR_Investor_Deck.pptx"
OUTPUT = "/Users/mohitmendiratta/Projects/agents/job1/MYTAR_Investor_Deck.pptx"


def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def rect_rnd(slide, left, top, width, height, fill_color, rad=0.05):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = rad
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def tb(slide, left, top, width, height, text, fs=None, color=None, bold=None,
       align=None, fname=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = fname
    if fs:   p.font.size = fs
    if color: p.font.color.rgb = color
    if bold is not None: p.font.bold = bold
    if align: p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    return txBox


def page_num(slide, num, total):
    tb(slide, PAGE_NUM_X, PAGE_NUM_Y, PAGE_NUM_W, PAGE_NUM_H,
       f"{num:02d} / {total:02d}", color=C_DIMGREY, align=PP_ALIGN.RIGHT)


def slide_base(slide, accent_color, label, title, subtitle):
    rect(slide, 0, 0, W, TOP_BAR_H, fill_color=accent_color)
    rect(slide, W - RIGHT_BAR_W, 0, RIGHT_BAR_W, H, fill_color=accent_color)
    tb(slide, LABEL_X, LABEL_Y, LABEL_W, LABEL_H, label, color=accent_color, bold=True)
    tb(slide, TITLE_X, TITLE_Y, TITLE_W, TITLE_H, title, fs=Pt(32), color=C_WHITE, bold=True)
    if subtitle:
        tb(slide, SUBTITLE_X, SUBTITLE_Y, SUBTITLE_W, SUBTITLE_H,
           subtitle, fs=Pt(14), color=C_LTGREY)
    rect(slide, ACCNT_LINE_X, ACCNT_LINE_Y, ACCNT_LINE_W, ACCNT_LINE_H,
         fill_color=accent_color)


#
# ── Slide 24: Market Fragmentation ── matches slide 23 pattern
#
def s24(prs):
    s = add_blank_slide(prs)
    slide_base(s, C_ORANGE, "24  |  COMPETITIVE",
               "Market Fragmentation",
               "The market is highly fragmented across four disconnected categories")

    # 4 top-strip cards exactly like slide 23
    cats = [
        ("Identity\nSystems", "▸  Verification\n▸  Credentials\n▸  Anti-bot proof", C_TEAL),
        ("Content\nAuthenticity", "▸  Content provenance\n▸  AI watermarking\n▸  Metadata tracking", C_BLUE),
        ("AI Avatar\nPlatforms", "▸  Avatar generation\n▸  Video synthesis\n▸  Enterprise content", C_PURPLE),
        ("AI\nCompanions", "▸  Emotional AI\n▸  Roleplay chars\n▸  Personal assistants", C_GREEN),
    ]
    card_w = 2560320
    card_h = 2560320
    card_x0 = 731520
    card_y = 2103120
    card_gap = 3566160 - (731520 + 2560320)  # gap = 274320

    for i, (title, body, color) in enumerate(cats):
        cx = card_x0 + i * (card_w + card_gap)
        rect_rnd(s, cx, card_y, card_w, card_h, fill_color=C_CARD)
        rect(s, cx, card_y, card_w, ACCNT_STRIP_H, fill_color=color)
        # Title inside card
        tb(s, cx + 182880, card_y + 182880, card_w - 365760, 640080,
           title, fs=Pt(18), color=C_WHITE, bold=True)
        # Body text
        tb(s, cx + 182880, card_y + 914400, card_w - 365760, 1463040,
           body, fs=Pt(12), color=C_LTGREY)

    # Bottom highlight
    rect(s, 731520, 5029200, 10698480, 1097280, fill_color=C_BOTDARK)
    rect(s, 731520, 5029200, 63500, 1097280, fill_color=C_ORANGE)
    tb(s, 1371600, 5120640, 9601200, 914400,
       "Despite rapid innovation in digital identity, AI avatars, and content authenticity, "
       "today\u2019s ecosystem is split into isolated vertical stacks, "
       "each solving only a narrow part of the problem.",
       fs=Pt(14), color=C_LTGREY)

    page_num(s, 24, 36)


#
# ── Slide 25: Identity Systems (TEAL) ── matches slide 24 left-accent-bar cards
#
def s25(prs):
    s = add_blank_slide(prs)
    slide_base(s, C_TEAL, "25  |  COMPETITOR ANALYSIS",
               "Identity Systems",
               "These platforms verify human identity but do not extend to AI representation")

    comps = [
        ("World ID", "Tools for Humanity", [
            "Proof-of-human via iris scanning",
            "Anti-bot / Sybil resistance",
            "No AI persona or monetization",
        ]),
        ("Microsoft Entra Verified ID", "Enterprise verifiable credentials", [
            "Enterprise-grade verifiable credentials",
            "Deep identity & compliance integration",
            "Not designed for consumer AI",
        ]),
    ]

    card_w = 5212080
    card_h = 2560320
    card_gap = 635000
    total_w = 2 * card_w + card_gap
    left_start = (W - total_w) // 2
    card_y = 2286000

    for i, (name, subtitle, bullets) in enumerate(comps):
        cx = left_start + i * (card_w + card_gap)
        # Card body
        rect_rnd(s, cx, card_y, card_w, card_h, fill_color=C_CARD)
        # Left accent bar
        rect(s, cx, card_y, LEFT_ACCNT_W, card_h, fill_color=C_TEAL)
        # Name
        tb(s, cx + 274320, card_y + 274320, card_w - 548640, 365760,
           name, fs=Pt(20), color=C_WHITE, bold=True)
        # Subtitle
        tb(s, cx + 274320, card_y + 731520, card_w - 548640, 274320,
           subtitle, fs=Pt(11), color=C_DIMGREY)
        # Bullets
        for j, b in enumerate(bullets):
            tb(s, cx + 274320, card_y + 1188720 + j * 457200,
               card_w - 548640, 365760,
               f"▸  {b}", fs=Pt(12), color=C_LTGREY)

    page_num(s, 25, 36)


#
# ── Slide 26: Content Authenticity (BLUE) ── same pattern
#
def s26(prs):
    s = add_blank_slide(prs)
    slide_base(s, C_BLUE, "26  |  COMPETITOR ANALYSIS",
               "Content Authenticity",
               "These systems validate AI-generated or traced content, "
               "but lack ownership and identity")

    comps = [
        ("Adobe Content Credentials", "Cryptographic media provenance", [
            "Cryptographic metadata for provenance",
            "Industry standard for attribution",
            "No link to persistent identity",
        ]),
        ("Google SynthID", "Invisible AI content watermarking", [
            "Invisible AI content watermarking",
            "Scalable detection + attribution",
            "No ownership or economic enforcement",
        ]),
    ]

    card_w = 5212080
    card_h = 2560320
    card_gap = 635000
    total_w = 2 * card_w + card_gap
    left_start = (W - total_w) // 2
    card_y = 2286000

    for i, (name, subtitle, bullets) in enumerate(comps):
        cx = left_start + i * (card_w + card_gap)
        rect_rnd(s, cx, card_y, card_w, card_h, fill_color=C_CARD)
        rect(s, cx, card_y, LEFT_ACCNT_W, card_h, fill_color=C_BLUE)
        tb(s, cx + 274320, card_y + 274320, card_w - 548640, 365760,
           name, fs=Pt(20), color=C_WHITE, bold=True)
        tb(s, cx + 274320, card_y + 731520, card_w - 548640, 274320,
           subtitle, fs=Pt(11), color=C_DIMGREY)
        for j, b in enumerate(bullets):
            tb(s, cx + 274320, card_y + 1188720 + j * 457200,
               card_w - 548640, 365760,
               f"▸  {b}", fs=Pt(12), color=C_LTGREY)

    page_num(s, 26, 36)


#
# ── Slide 27: AI Avatars (PURPLE) ── 3 narrow left-accent-bar cards (like slide 9)
#
def s27(prs):
    s = add_blank_slide(prs)
    slide_base(s, C_PURPLE, "27  |  COMPETITOR ANALYSIS",
               "AI Avatar / Synthetic Identity Platforms",
               "These tools create digital representations but lack identity anchoring")

    comps = [
        ("HeyGen", [
            ("Strength", "High-quality avatar\nsynthesis for enterprise"),
            ("Limitation", "Avatars not persistent\nidentity-linked assets"),
        ]),
        ("Synthesia", [
            ("Strength", "Scalable enterprise\nvideo generation"),
            ("Limitation", "No persistent digital\nself layer"),
        ]),
        ("D-ID", [
            ("Strength", "Lightweight avatar gen\nfrom static images"),
            ("Limitation", "No identity binding,\ncommerce, or memory"),
        ]),
    ]

    card_w = 3383280
    card_h = 3474720
    card_gap = 50800
    total_w = 3 * card_w + 2 * card_gap
    left_start = (W - total_w) // 2
    card_y = 2103120

    for i, (name, details) in enumerate(comps):
        cx = left_start + i * (card_w + card_gap)
        rect_rnd(s, cx, card_y, card_w, card_h, fill_color=C_CARD)
        rect(s, cx, card_y, LEFT_ACCNT_W2, card_h, fill_color=C_PURPLE)
        # Name centered
        tb(s, cx + 182880, card_y + 274320, card_w - 365760, 365760,
           name, fs=Pt(18), color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
        # Details
        for j, (label, value) in enumerate(details):
            dy = card_y + 822960 + j * 1188720
            tb(s, cx + 182880, dy, card_w - 365760, 228600,
               label, fs=Pt(10), color=C_PURPLE, bold=True, align=PP_ALIGN.CENTER)
            tb(s, cx + 182880, dy + 274320, card_w - 365760, 731520,
               value, fs=Pt(11), color=C_LTGREY, align=PP_ALIGN.CENTER)

    page_num(s, 27, 36)


#
# ── Slide 28: AI Companions (GREEN) ── same 2-column pattern as S25
#
def s28(prs):
    s = add_blank_slide(prs)
    slide_base(s, C_GREEN, "28  |  COMPETITOR ANALYSIS",
               "AI Companions",
               "These platforms focus on conversational personality and engagement")

    comps = [
        ("Replika", "Emotional AI companionship", [
            "Deep personalization & memory",
            "Relationship simulation",
            "No verified identity or commerce",
        ]),
        ("Character AI", "User-generated AI characters", [
            "Massive creative character ecosystem",
            "Rich roleplay experiences",
            "No persistent identity layer",
        ]),
    ]

    card_w = 5212080
    card_h = 2560320
    card_gap = 635000
    total_w = 2 * card_w + card_gap
    left_start = (W - total_w) // 2
    card_y = 2286000

    for i, (name, subtitle, bullets) in enumerate(comps):
        cx = left_start + i * (card_w + card_gap)
        rect_rnd(s, cx, card_y, card_w, card_h, fill_color=C_CARD)
        rect(s, cx, card_y, LEFT_ACCNT_W, card_h, fill_color=C_GREEN)
        tb(s, cx + 274320, card_y + 274320, card_w - 548640, 365760,
           name, fs=Pt(20), color=C_WHITE, bold=True)
        tb(s, cx + 274320, card_y + 731520, card_w - 548640, 274320,
           subtitle, fs=Pt(11), color=C_DIMGREY)
        for j, b in enumerate(bullets):
            tb(s, cx + 274320, card_y + 1188720 + j * 457200,
               card_w - 548640, 365760,
               f"▸  {b}", fs=Pt(12), color=C_LTGREY)

    page_num(s, 28, 36)


#
# ── Slide 29: Core Insight (ORANGE) ── summary
#
def s29(prs):
    s = add_blank_slide(prs)
    slide_base(s, C_ORANGE, "29  |  KEY INSIGHT",
               "Core Insight \u2014 The Missing Layer",
               "Across all categories, systems optimize for one dimension only")

    # 4 compact cards showing what each does
    dims = [
        ("\U0001F511", "Identity", "\u2192 verify who\nyou are", C_TEAL),
        ("\u2705", "Authenticity", "\u2192 verify what\nis real", C_BLUE),
        ("\U0001F5E3\ufe0f", "Avatars", "\u2192 generate what\nyou look like", C_PURPLE),
        ("\U0001F9D1\u200D\U0001F4BB", "Companions", "\u2192 simulate how\nyou talk", C_GREEN),
    ]

    card_w = 2560320
    card_h = 2194560
    card_x0 = 731520
    card_y = 2194560
    card_gap = 274320

    for i, (icon, title, desc, color) in enumerate(dims):
        cx = card_x0 + i * (card_w + card_gap)
        rect_rnd(s, cx, card_y, card_w, card_h, fill_color=C_CARD)
        rect(s, cx, card_y, card_w, ACCNT_STRIP_H, fill_color=color)
        tb(s, cx + 182880, card_y + 274320, card_w - 365760, 457200,
           icon, fs=Pt(26), color=color, align=PP_ALIGN.CENTER)
        tb(s, cx + 182880, card_y + 731520, card_w - 365760, 274320,
           title, fs=Pt(12), color=color, bold=True, align=PP_ALIGN.CENTER)
        tb(s, cx + 182880, card_y + 1097280, card_w - 365760, 914400,
           desc, fs=Pt(14), color=C_LTGREY, align=PP_ALIGN.CENTER)

    # Bottom bar with insight
    bottom_y = 4754880
    rect(s, 731520, bottom_y, 10698480, 914400, fill_color=C_BOTDARK)
    rect(s, 731520, bottom_y, 63500, 914400, fill_color=C_ORANGE)
    tb(s, 914400, bottom_y + 91440, 5029200, 365760,
       "But none combine the full stack:",
       fs=Pt(16), color=C_WHITE, bold=True)
    tb(s, 914400, bottom_y + 457200, 10058400, 365760,
       "Identity + Ownership + AI Twin + Content Provenance + Commerce",
       fs=Pt(18), color=C_ORANGE, bold=True, align=PP_ALIGN.CENTER)

    page_num(s, 29, 36)


#
# ── Slide 30: Strategic Implication (ORANGE)
#
def s30(prs):
    s = add_blank_slide(prs)
    slide_base(s, C_ORANGE, "30  |  STRATEGIC IMPLICATION",
               "Strategic Implication for MYTAR",
               "This fragmentation creates a structural gap in the market")

    # Two large cards side by side
    card_w = 5029200
    card_h = 3840480
    card_gap = 635000
    total_w = 2 * card_w + card_gap
    left_start = (W - total_w) // 2
    card_y = 2194560

    # Left: Fragmentation Gap
    cx = left_start
    rect_rnd(s, cx, card_y, card_w, card_h, fill_color=C_CARD)
    tb(s, cx + 182880, card_y + 274320, card_w - 365760, 365760,
       "The Fragmentation Gap", fs=Pt(16), color=C_RED, bold=True)
    gaps = [
        "No persistent \u201cdigital self\u201d layer",
        "No unified AI-native identity graph",
        "No ownership for AI actions/content",
        "No monetization for AI twins",
    ]
    for j, g in enumerate(gaps):
        tb(s, cx + 182880, card_y + 822960 + j * 502920,
           card_w - 365760, 457200,
           f"\u2716  {g}", fs=Pt(13), color=C_LTGREY)

    # Right: Whitespace
    cx = left_start + card_w + card_gap
    rect_rnd(s, cx, card_y, card_w, card_h, fill_color=C_BOTTEAL)
    tb(s, cx + 182880, card_y + 274320, card_w - 365760, 365760,
       "The Whitespace MYTAR Targets", fs=Pt(16), color=C_TEAL, bold=True)
    ws = [
        "Identity is verifiable",
        "AI twins are persistent",
        "Content is attributable",
        "Actions are ownable",
        "Commerce is native to the\ndigital self",
    ]
    for j, w in enumerate(ws):
        tb(s, cx + 182880, card_y + 822960 + j * 502920,
           card_w - 365760, 457200,
           f"\u2713  {w}", fs=Pt(13), color=C_TEAL, bold=True)

    # Bottom caption
    tb(s, 731520, 6583680, 10058400, 365760,
       "This is the structural whitespace MYTAR targets: A unified system for the AI era.",
       fs=Pt(13), color=C_DIMGREY, align=PP_ALIGN.CENTER)

    page_num(s, 30, 36)


#
# ── Reorder slides ──
#
def reorder_slides(prs, new_order):
    ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    pres_elem = prs.part._element
    sldIdLst = pres_elem.find(f'{{{ns}}}sldIdLst')
    sldId_elems = list(sldIdLst)
    ordered = [sldId_elems[i] for i in new_order]
    for elem in sldId_elems:
        sldIdLst.remove(elem)
    for elem in ordered:
        sldIdLst.append(elem)


#
# ── Main ──
#
def main():
    prs = Presentation(INPUT)
    orig_count = len(prs.slides)
    print(f"Original slides: {orig_count}")

    # Create new competition slides (appended at end)
    print("Creating competition slides (v1 style)...")
    s24(prs)
    s25(prs)
    s26(prs)
    s27(prs)
    s28(prs)
    s29(prs)
    s30(prs)

    new_count = len(prs.slides)
    print(f"Total slides after append: {new_count}")

    # New order: 0-22 (orig 1-23), then new slides (appended at end), then orig 24-29
    new_slide_start = orig_count  # where the 7 new slides start (0-indexed)
    new_order = list(range(23))  # 0-22
    new_order.extend(range(new_slide_start, new_slide_start + 7))  # new slides
    new_order.extend(range(23, 29))  # orig slides 23-28

    print("Reordering...")
    reorder_slides(prs, new_order)

    # Update page numbers on all slides
    print("Updating page numbers...")
    total = new_count
    for i, slide in enumerate(prs.slides):
        snum = i + 1
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if '/' in text and len(text) < 12:
                    tf = shape.text_frame
                    for para in tf.paragraphs:
                        if '/' in para.text:
                            para.text = f"{snum:02d} / {total:02d}"
                            break
                    break

    # Update section labels on moved slides (now at positions 31-36, 0-indexed 30-35)
    print("Updating section labels on moved slides...")
    label_map = {30: 31, 31: 32, 32: 33, 33: 34, 34: 35, 35: 36}
    for idx, new_num in label_map.items():
        slide = prs.slides[idx]
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                for para in tf.paragraphs:
                    old = para.text.strip()
                    m = __import__('re').match(r'^(\d{1,2})(\s*\|\s*.*)$', old)
                    if m:
                        para.text = f'{new_num:02d}{m.group(2)}'
                        break
                    m2 = __import__('re').match(r'^(\d{1,2})$', old)
                    if m2:
                        para.text = f'{new_num:02d}'
                        break

    prs.save(OUTPUT)
    print(f"Saved! Total slides: {len(prs.slides)}")
    print(f"File: {OUTPUT}")


if __name__ == "__main__":
    main()
